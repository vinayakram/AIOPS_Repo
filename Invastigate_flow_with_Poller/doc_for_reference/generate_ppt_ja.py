"""
Generate Investigation Pipeline Agent Documentation PPT (Japanese)
Run: python doc_for_reference/generate_ppt_ja.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ─────────────────────────────────────────────────────────
C_NAVY       = RGBColor(0x0D, 0x1B, 0x2A)
C_DARK_BLUE  = RGBColor(0x1A, 0x35, 0x5E)
C_ACCENT     = RGBColor(0x2E, 0x86, 0xAB)
C_GREEN      = RGBColor(0x27, 0xAE, 0x60)
C_AMBER      = RGBColor(0xF3, 0x9C, 0x12)
C_RED        = RGBColor(0xC0, 0x39, 0x2B)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GREY = RGBColor(0xCC, 0xD6, 0xE0)
C_BOX_BG     = RGBColor(0x16, 0x27, 0x3E)
C_BOX_BG2    = RGBColor(0x1E, 0x34, 0x52)
C_STEP_BG    = RGBColor(0x0A, 0x42, 0x6B)


# ── Helpers ────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=12, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        5,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def slide_header(slide, title, subtitle=None, accent_bar=True):
    if accent_bar:
        bar = add_rect(slide, 0, 0, 13.33, 0.75, C_DARK_BLUE)
        add_textbox(slide, 0.3, 0.1, 12.5, 0.55, title,
                    font_size=20, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, 0.3, 0.78, 12.5, 0.35, subtitle,
                    font_size=11, color=C_LIGHT_GREY, italic=True)


def bullet_block(slide, left, top, width, height, title, bullets,
                 title_color=C_ACCENT, bg=C_BOX_BG, bullet_color=C_WHITE,
                 title_size=13, bullet_size=10.5):
    add_rect(slide, left, top, width, height, bg)
    add_textbox(slide, left + 0.15, top + 0.05, width - 0.2, 0.3,
                title, font_size=title_size, bold=True, color=title_color)
    y = top + 0.35
    for b in bullets:
        add_textbox(slide, left + 0.25, y, width - 0.35, 0.28,
                    f"• {b}", font_size=bullet_size, color=bullet_color)
        y += 0.28
    return y


def step_box(slide, left, top, width, num, title, desc, num_color=C_AMBER):
    add_rounded_rect(slide, left, top, width, 0.72, C_STEP_BG, C_ACCENT)
    add_textbox(slide, left + 0.08, top + 0.06, 0.35, 0.55,
                str(num), font_size=22, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.48, top + 0.06, width - 0.58, 0.28,
                title, font_size=11, bold=True, color=C_WHITE)
    add_textbox(slide, left + 0.48, top + 0.33, width - 0.58, 0.35,
                desc, font_size=9.5, color=C_LIGHT_GREY)


def table_row(slide, left, top, cols, values, colors, heights=0.32, bold_flags=None):
    x = left
    for i, (col_w, val) in enumerate(zip(cols, values)):
        add_rect(slide, x, top, col_w, heights, colors[i])
        bf = bold_flags[i] if bold_flags else False
        add_textbox(slide, x + 0.08, top + 0.04, col_w - 0.1, heights - 0.05,
                    str(val), font_size=9.5, bold=bf, color=C_WHITE)
        x += col_w


def arrow_right(slide, left, top, width=0.4):
    add_textbox(slide, left, top, width, 0.25, "→",
                font_size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def arrow_down(slide, left, top, height=0.25):
    add_textbox(slide, left, top, 0.35, height, "▼",
                font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ── Presentation setup ─────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════════════════
# スライド 1 — タイトル
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)

add_rect(slide, 0, 2.8, 13.33, 0.08, C_ACCENT)

add_textbox(slide, 1.0, 0.9, 11.33, 1.2,
            "調査パイプライン",
            font_size=42, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 1.0, 2.05, 11.33, 0.7,
            "マルチエージェント可観測性システム — 技術ドキュメント",
            font_size=20, color=C_ACCENT, align=PP_ALIGN.CENTER)

agents = [
    ("1", "正規化"), ("2", "相関分析"), ("3", "エラー分析"),
    ("4", "根本原因分析"), ("5", "推奨"),
]
x = 1.0
for num, name in agents:
    add_rounded_rect(slide, x, 3.2, 2.1, 0.9, C_DARK_BLUE, C_ACCENT)
    add_textbox(slide, x, 3.22, 2.1, 0.38, num,
                font_size=16, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, 3.58, 2.1, 0.42, name,
                font_size=11, color=C_WHITE, align=PP_ALIGN.CENTER)
    if num != "5":
        add_textbox(slide, x + 2.1, 3.5, 0.2, 0.3, "→",
                    font_size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    x += 2.27

add_textbox(slide, 1.0, 4.35, 11.33, 0.4,
            "Langfuse  ·  Prometheus  ·  GPT-4o  ·  Pydantic  ·  FastAPI",
            font_size=13, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 6.9, 12.33, 0.35,
            "Prodapt  ·  AI可観測性チーム  ·  2026",
            font_size=10, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════
# スライド 2 — アジェンダ
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "アジェンダ")

items = [
    ("01", "システム概要",              "アーキテクチャ、主要原則、技術スタック"),
    ("02", "データソース",              "Langfuse（AIトレース）とPrometheus（インフラメトリクス）"),
    ("03", "エージェント1 — 正規化",    "生ログ取り込み、エラー分類、NO_ERROR短絡処理"),
    ("04", "エージェント2 — 相関分析",  "クロスシステム因果グラフ、analysis_targetルーティング"),
    ("05", "エージェント3 — エラー分析","詳細エラー抽出、パターン検出、エラーID付与"),
    ("06", "エージェント4 — 根本原因分析","根本原因特定、因果連鎖、5つのなぜ分析"),
    ("07", "エージェント5 — 推奨",      "ランク付けされた実行可能ソリューション、工数見積"),
    ("08", "パイプラインデータフロー",  "エージェント間のエンドツーエンドデータ受け渡し"),
    ("09", "効率性・設計方針",          "スキーマ契約、短絡処理、グレースフルデグラデーション"),
]

y = 1.05
for num, title, desc in items:
    add_rect(slide, 0.4, y, 0.5, 0.52, C_ACCENT)
    add_textbox(slide, 0.4, y + 0.08, 0.5, 0.35, num,
                font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.9, y, 11.9, 0.52, C_BOX_BG)
    add_textbox(slide, 1.05, y + 0.04, 3.8, 0.26, title,
                font_size=12, bold=True, color=C_WHITE)
    add_textbox(slide, 1.05, y + 0.28, 11.0, 0.22, desc,
                font_size=9.5, color=C_LIGHT_GREY)
    y += 0.6


# ══════════════════════════════════════════════════════════════════════════
# スライド 3 — システム概要
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "システム概要", "順次実行の5エージェントパイプライン — 各エージェントは明確に範囲が定義された1つの責務を担う")

agents_flow = [
    ("1", "正規化",       "エラー種別分類\n信号抽出", C_BOX_BG),
    ("2", "相関分析",     "因果グラフ構築\n分析ルーティング", C_BOX_BG),
    ("3", "エラー分析",   "エラー抽出・ID付与\nパターン検出", C_BOX_BG),
    ("4", "根本原因分析", "根本原因特定\n5つのなぜ分析", C_BOX_BG),
    ("5", "推奨",         "ランク付き\n実行可能ソリューション（1〜4件）", C_BOX_BG),
]

x = 0.3
for num, name, desc, bg in agents_flow:
    add_rounded_rect(slide, x, 1.1, 2.35, 1.6, bg, C_ACCENT)
    add_rect(slide, x, 1.1, 2.35, 0.38, C_ACCENT)
    add_textbox(slide, x, 1.12, 2.35, 0.35,
                f"ステップ{num}  •  {name}", font_size=11, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 1.55, 2.15, 0.9, desc,
                font_size=9.5, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER)
    if num != "5":
        add_textbox(slide, x + 2.35, 1.72, 0.25, 0.3, "→",
                    font_size=18, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    x += 2.6

principles = [
    ("決定論的", "全LLM呼び出しにtemperature=0.0 — 再現可能な出力"),
    ("スキーマ駆動", "Pydanticモデルが各プロンプトに注入するJSONスキーマを自動生成"),
    ("グレースフル", "全ての外部障害はWARNプレースホルダーを生成 — エージェントはクラッシュしない"),
    ("分離型", "各エージェントは直前のエージェント出力のみ参照 — 共有状態なし"),
]
y = 3.0
add_textbox(slide, 0.3, 2.9, 12.7, 0.3, "主要原則",
            font_size=10, bold=True, color=C_ACCENT)
for i, (title, desc) in enumerate(principles):
    col = i % 2
    row = i // 2
    lx = 0.3 + col * 6.55
    ly = y + 0.05 + row * 0.95
    add_rect(slide, lx, ly, 6.3, 0.82, C_BOX_BG2)
    add_rect(slide, lx, ly, 0.08, 0.82, C_GREEN)
    add_textbox(slide, lx + 0.2, ly + 0.06, 6.0, 0.28,
                title, font_size=11, bold=True, color=C_GREEN)
    add_textbox(slide, lx + 0.2, ly + 0.36, 6.0, 0.38,
                desc, font_size=9.5, color=C_LIGHT_GREY)

add_textbox(slide, 0.3, 6.85, 12.5, 0.28,
            "技術スタック:  GPT-4o  ·  FastAPI  ·  Pydantic v2  ·  AsyncOpenAI  ·  httpx  ·  SQLite",
            font_size=9, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════
# スライド 4 — データソース: Langfuse
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "データソース1 — Langfuse（AIエージェントトレース）")

add_rect(slide, 0.3, 1.0, 5.9, 5.8, C_BOX_BG)
add_textbox(slide, 0.5, 1.05, 5.6, 0.3, "概要",
            font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 0.5, 1.35, 5.6, 0.55,
            "LLM可観測性プラットフォーム。AIエージェントの全インタラクションをトレース・スパン・ジェネレーションとして記録する",
            font_size=10, color=C_WHITE)

add_textbox(slide, 0.5, 1.95, 5.6, 0.3, "認証方式",
            font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 0.5, 2.25, 5.6, 0.3,
            "HTTP Basic認証  →  public_key : secret_key",
            font_size=10, color=C_LIGHT_GREY)

add_textbox(slide, 0.5, 2.65, 5.6, 0.3, "1回のフェッチで発行するAPIコール（計2件）",
            font_size=10, bold=True, color=C_ACCENT)

calls = [
    ("GET", "/api/public/traces/{trace_id}", "トップレベルトレース: 名称・ステータス・入出力・レイテンシ・コスト"),
    ("GET", "/api/public/observations?traceId=...&limit=100", "子スパン/ジェネレーション全件: ステータス・モデル・トークン数・エラー"),
]
y = 2.98
for method, endpoint, desc in calls:
    add_rect(slide, 0.5, y, 0.55, 0.28, C_GREEN)
    add_textbox(slide, 0.5, y + 0.02, 0.55, 0.24, method,
                font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.1, y, 4.7, 0.28, endpoint,
                font_size=9, color=C_AMBER)
    add_textbox(slide, 0.5, y + 0.3, 5.6, 0.28, desc,
                font_size=9, color=C_LIGHT_GREY)
    y += 0.68

add_textbox(slide, 0.5, 4.4, 5.6, 0.3, "使用エージェント",
            font_size=10, bold=True, color=C_ACCENT)
agents_using = ["正規化（trace_id指定時）", "相関分析（trace_id利用可能時は常時）",
                "エラー分析（対象=AgentまたはUnknown時）", "根本原因分析（対象=AgentまたはUnknown時）"]
yy = 4.7
for a in agents_using:
    add_textbox(slide, 0.6, yy, 5.4, 0.24, f"• {a}", font_size=9.5, color=C_WHITE)
    yy += 0.26

add_rect(slide, 6.5, 1.0, 6.55, 5.8, C_BOX_BG2)
add_textbox(slide, 6.7, 1.05, 6.2, 0.3, "エラー検出ロジック",
            font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 6.7, 1.35, 6.2, 0.35,
            "各observationを4フィールドの順序で ERROR / WARN / INFO に分類:",
            font_size=9.5, color=C_WHITE)

detection_steps = [
    ("1", "statusフィールド",   "ERROR / FAIL / FAILED → ERROR"),
    ("2", "statusMessage",      "エラーキーワード含有 → ERROR"),
    ("3", "outputフィールド",   "JSONに'error'キー → ERROR  |  文字列マッチ → ERROR"),
    ("4", "inputフィールド",    "エラーキーワードマッチ → WARN"),
    ("5", "デフォルト",         "→ INFO"),
]
y = 1.8
for num, field, rule in detection_steps:
    add_rect(slide, 6.6, y, 0.35, 0.38, C_ACCENT)
    add_textbox(slide, 6.6, y + 0.04, 0.35, 0.3, num,
                font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 7.0, y + 0.02, 1.5, 0.18, field,
                font_size=9.5, bold=True, color=C_AMBER)
    add_textbox(slide, 7.0, y + 0.2, 5.8, 0.2, rule,
                font_size=9, color=C_LIGHT_GREY)
    y += 0.5

add_textbox(slide, 6.7, 4.55, 6.2, 0.3, "出力フォーマット（ログエントリ単位）",
            font_size=10, bold=True, color=C_ACCENT)
fields = ["timestamp（タイムスタンプ）", "source  (langfuse)", "service  （observation名）",
          "message  （種別 + ステータス + エラー + トークン数）", "level  (ERROR / WARN / INFO)",
          "metadata  （各種ID・モデル・コスト・レイテンシ・入出力）"]
y = 4.85
for f in fields:
    add_textbox(slide, 6.8, y, 6.1, 0.24, f"→  {f}", font_size=9.5, color=C_WHITE)
    y += 0.26


# ══════════════════════════════════════════════════════════════════════════
# スライド 5 — データソース: Prometheus
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "データソース2 — Prometheus（インフラメトリクス）")

add_rect(slide, 0.3, 1.0, 12.7, 1.1, C_BOX_BG)
add_textbox(slide, 0.5, 1.05, 8.0, 0.3, "クエリ方式",
            font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 0.5, 1.35, 12.2, 0.55,
            "GET /api/v1/query_range  ·  {agent}テンプレートのPromQL  ·  インシデントタイムスタンプ前後±5分  ·  15秒ステップ分解能  ·  HTTP 10秒タイムアウト",
            font_size=10, color=C_WHITE)

add_textbox(slide, 0.3, 2.2, 12.7, 0.3, "デフォルト6件のPromQLクエリ（フェッチ毎に実行）",
            font_size=10, bold=True, color=C_ACCENT)

cols = [2.2, 5.5, 2.5, 2.3]
table_row(slide, 0.3, 2.55, cols,
          ["クエリ名", "PromQLパターン", "検出対象", "重大度ルール"],
          [C_DARK_BLUE, C_DARK_BLUE, C_DARK_BLUE, C_DARK_BLUE],
          heights=0.32, bold_flags=[True, True, True, True])

rows = [
    ("error_rate",    "rate(http_requests_total{status=~'5..',job=~'.*{agent}.*'}[window])", "HTTP 5xxエラー率",      "ERROR: > 0"),
    ("latency_p99",   "histogram_quantile(0.99, rate(http_request_duration_seconds_bucket...))",  "P99レイテンシスパイク", "INFO のみ"),
    ("up_status",     "up{job=~'.*{agent}.*'}",                                              "サービス稼働/停止",       "ERROR: == 0"),
    ("memory_usage",  "container_memory_usage_bytes{pod=~'.*{agent}.*'}",                    "コンテナメモリ使用量",     "INFO のみ"),
    ("restart_count", "kube_pod_container_status_restarts_total{pod=~'.*{agent}.*'}",         "Podクラッシュループ検出",  "WARN: > 0"),
    ("dns_failures",  "rate(coredns_dns_responses_total{rcode='SERVFAIL'}[window])",          "DNS SERVFAIL率",          "ERROR: > 0"),
]
y = 2.87
for i, (name, promql, detects, severity) in enumerate(rows):
    bg = C_BOX_BG if i % 2 == 0 else C_BOX_BG2
    sev_color = C_RED if "ERROR" in severity else C_AMBER if "WARN" in severity else C_BOX_BG
    table_row(slide, 0.3, y, cols,
              [name, promql, detects, severity],
              [bg, bg, bg, sev_color], heights=0.3)
    y += 0.3

add_rect(slide, 0.3, 5.6, 4.0, 1.6, C_BOX_BG)
add_textbox(slide, 0.5, 5.65, 3.7, 0.28, "時間ウィンドウ", font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 0.5, 5.95, 3.7, 0.9,
            "インシデントタイムスタンプ ± 5分\n障害の前後両方をキャプチャ",
            font_size=10, color=C_WHITE)

add_rect(slide, 4.5, 5.6, 4.2, 1.6, C_BOX_BG)
add_textbox(slide, 4.7, 5.65, 3.9, 0.28, "クエリ失敗時の処理", font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 4.7, 5.95, 3.9, 0.9,
            "各クエリ失敗はWARNプレースホルダーエントリを生成。LLMにデータ欠損を通知 — パイプラインは中断しない",
            font_size=10, color=C_WHITE)

add_rect(slide, 8.9, 5.6, 4.1, 1.6, C_BOX_BG)
add_textbox(slide, 9.1, 5.65, 3.8, 0.28, "データポイント選択", font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 9.1, 5.95, 3.8, 0.9,
            "時系列ごとに最後のデータポイントのみ使用 — インシデント発生時点の最新状態",
            font_size=10, color=C_WHITE)


# ══════════════════════════════════════════════════════════════════════════
# スライド 6 — エージェント1: 正規化
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "エージェント1 — 正規化エージェント",
             "目的: 非構造化の生ログを単一の構造化NormalizedIncidentオブジェクトに変換する")

add_rect(slide, 0.3, 1.0, 8.6, 0.55, C_BOX_BG)
add_rect(slide, 0.3, 1.0, 0.08, 0.55, C_AMBER)
add_textbox(slide, 0.5, 1.06, 8.3, 0.42,
            "厳格なスコープ: エラー種別の分類と原子的信号の抽出のみ。因果関係・根本原因・相関分析は行わない。",
            font_size=10, bold=True, color=C_AMBER)

add_textbox(slide, 0.3, 1.7, 12.7, 0.28, "内部処理ステップ", font_size=10, bold=True, color=C_ACCENT)

steps = [
    ("1", "データソース振り分け",     "trace_id あり → Langfuse    |    trace_id なし → Prometheus"),
    ("2", "生ログ取得",               "Langfuse: 1トレース + 最大100 observations    |    Prometheus: 6件のPromQL範囲クエリ"),
    ("3", "LLM前エラースキャン",      "_has_error_signals(): レベルチェック（ERROR/WARN） + 完全単語キーワード正規表現スキャン"),
    ("4", "NO_ERROR短絡処理",         "信号未検出 → 即時NO_ERROR返却、confidence=1.0 — LLMは呼び出さない"),
    ("5", "LLMプロンプト構築",        "ログを [N] ts=... svc=... level=... msg=... meta={...} 形式でフォーマット + Pydanticスキーマ注入"),
    ("6", "GPT-4o呼び出し",           "temperature=0.0  ·  response_format=json_object  ·  スキーマをプロンプトで強制"),
    ("7", "Pydanticバリデーション",   "NormalizedIncident.model_validate(data) — スキーマ不一致は返却前に拒否"),
]
y = 2.05
for num, title, desc in steps:
    step_box(slide, 0.3, y, 12.7, num, title, desc)
    y += 0.78

add_rect(slide, 9.15, 1.0, 4.15, 0.55, C_BOX_BG2)
add_textbox(slide, 9.35, 1.06, 3.9, 0.42, "ルーティング:  trace_id → Langfuse   |   trace_idなし → Prometheus",
            font_size=9.5, color=C_LIGHT_GREY, italic=True)


# ══════════════════════════════════════════════════════════════════════════
# スライド 7 — エージェント1: 出力契約・効率性
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "エージェント1 — 出力契約と効率性機能")

add_textbox(slide, 0.3, 1.0, 6.0, 0.3, "出力: NormalizedIncident", font_size=11, bold=True, color=C_ACCENT)
fields = [
    ("error_type",    "NO_ERROR | AI_AGENT | INFRA | NETWORK | UNKNOWN",  "優先度: INFRA > NETWORK > AI_AGENT（複合時）"),
    ("error_summary", "str（最大300文字）",                                "事実に基づく1〜2行の説明 — 推測なし"),
    ("timestamp",     "ISO-8601文字列",                                    "ログから抽出した最初の発生時刻"),
    ("confidence",    "float  0.0 – 1.0",                                  "ログ証拠の明示度に基づく"),
    ("entities",      "{ agent_id, service, trace_id }",                   "抽出されたエンティティ参照"),
    ("signals",       "list[str]",                                         "例: ['LLM_access_disabled', 'timeout']"),
]
y = 1.35
for fname, ftype, fdesc in fields:
    add_rect(slide, 0.3, y, 6.1, 0.62, C_BOX_BG)
    add_textbox(slide, 0.45, y + 0.04, 2.0, 0.24, fname, font_size=10, bold=True, color=C_AMBER)
    add_textbox(slide, 2.5, y + 0.04, 3.8, 0.24, ftype, font_size=10, color=C_WHITE)
    add_textbox(slide, 0.45, y + 0.32, 5.8, 0.24, fdesc, font_size=9, color=C_LIGHT_GREY)
    y += 0.68

add_textbox(slide, 6.7, 1.0, 6.3, 0.3, "主要な効率性機能", font_size=11, bold=True, color=C_ACCENT)

features = [
    ("NO_ERROR短絡処理",
     "パイプライン全体で最もインパクトのある最適化。_has_error_signals()がFalseを返すとLLMは一切呼び出されず、パイプラインが終了する。本番環境で最も頻繁なケース（正常トレース）であり、実行ごとに約$0.10〜$0.30を節約。"),
    ("2フェーズエラー検出",
     "フェーズ1: 明示的レベルチェック（O(1)辞書参照）。フェーズ2: 曖昧レベルのみキーワードスキャン。'error_rate=0'などのメトリクス名による誤検知を防ぐ、事前コンパイル済み完全単語正規表現を使用。"),
    ("起動時スキーマ構築・キャッシュ",
     "NormalizedIncident.model_json_schema()は__init__で1回だけ呼び出され、self._response_schemaとしてキャッシュ。全リクエストが同一スキーマ文字列を再利用 — リフレクションの繰り返し処理なし。"),
    ("非同期ノンブロッキング",
     "AsyncOpenAI + httpx.AsyncClient — 全I/Oはノンブロッキング。FastAPIがスレッドなしで並行リクエストを処理。"),
]

y = 1.35
for title, desc in features:
    add_rect(slide, 6.7, y, 6.3, 1.32, C_BOX_BG2)
    add_rect(slide, 6.7, y, 0.07, 1.32, C_GREEN)
    add_textbox(slide, 6.85, y + 0.08, 6.1, 0.28, title, font_size=10.5, bold=True, color=C_GREEN)
    add_textbox(slide, 6.85, y + 0.4, 6.1, 0.82, desc, font_size=9.5, color=C_WHITE)
    y += 1.42


# ══════════════════════════════════════════════════════════════════════════
# スライド 8 — エージェント2: 相関分析
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "エージェント2 — 相関分析エージェント",
             "目的: クロスシステムの因果障害グラフを構築し、以降の分析ルーティング先を決定する")

add_rect(slide, 0.3, 1.0, 12.7, 0.55, C_BOX_BG)
add_rect(slide, 0.3, 1.0, 0.08, 0.55, C_AMBER)
add_textbox(slide, 0.5, 1.06, 12.3, 0.42,
            "重要な出力: analysis_target  —  以降の全エージェントが使用するデータソースを制御するルーティング決定",
            font_size=10, bold=True, color=C_AMBER)

add_textbox(slide, 0.3, 1.7, 6.0, 0.28, "データ取得戦略", font_size=10, bold=True, color=C_ACCENT)

add_rect(slide, 0.3, 2.02, 5.9, 0.58, C_GREEN)
add_textbox(slide, 0.5, 2.09, 5.6, 0.22, "常時取得", font_size=9.5, bold=True, color=C_NAVY)
add_textbox(slide, 0.5, 2.32, 5.6, 0.24, "Prometheus  →  6件のPromQLクエリ  |  インフラベースライン", font_size=9.5, color=C_NAVY)

add_rect(slide, 0.3, 2.65, 5.9, 0.58, C_ACCENT)
add_textbox(slide, 0.5, 2.72, 5.6, 0.22, "条件付き（trace_id利用可能時）", font_size=9.5, bold=True, color=C_WHITE)
add_textbox(slide, 0.5, 2.95, 5.6, 0.24, "Langfuse  →  トレース + 最大100 observations  |  AIレイヤー視点", font_size=9.5, color=C_WHITE)

add_rect(slide, 0.3, 3.28, 5.9, 0.45, C_BOX_BG)
add_textbox(slide, 0.5, 3.35, 5.6, 0.35,
            "LLM前に全ログを時系列ソート — 因果関係の特定に必須（最初の障害=根本候補）",
            font_size=9.5, color=C_LIGHT_GREY, italic=True)

add_textbox(slide, 0.3, 3.85, 6.0, 0.28, "処理ステップ", font_size=10, bold=True, color=C_ACCENT)
steps = [
    "Prometheus（常時）+ Langfuse（trace_id存在時）を取得",
    "両ソースの全ログを時系列ソート",
    "ユーザーメッセージ構築: ソース別ログ + 正規化コンテキスト",
    "CorrelationResultスキーマでGPT-4o呼び出し",
    "Pydanticバリデーションと CorrelationResult + analysis_target の返却",
]
y = 4.18
for i, s in enumerate(steps):
    add_rect(slide, 0.3, y, 0.35, 0.3, C_ACCENT)
    add_textbox(slide, 0.3, y + 0.03, 0.35, 0.24, str(i+1),
                font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.7, y + 0.03, 5.4, 0.28, s, font_size=9.5, color=C_WHITE)
    y += 0.38

add_textbox(slide, 6.7, 1.7, 6.3, 0.28, "出力: CorrelationResult", font_size=10, bold=True, color=C_ACCENT)

output_fields = [
    ("correlation_chain",    "list[str]  —  例: ['DNS障害 → プロキシタイムアウト → ゲートウェイ502']"),
    ("peer_components",      "list[PeerComponent]  —  役割+証拠を持つ各コンポーネント"),
    ("timeline",             "list[TimelineEvent]  —  時系列障害イベント"),
    ("root_cause_candidate", "コンポーネント + 信頼度(0〜1) + 理由"),
    ("analysis_target",      "Agent | InfraLogs | Unknown  ←  ルーティング決定"),
]
y = 2.02
for fname, fdesc in output_fields:
    bg = C_RED if fname == "analysis_target" else C_BOX_BG
    add_rect(slide, 6.7, y, 6.3, 0.52, bg)
    add_textbox(slide, 6.85, y + 0.04, 2.2, 0.22, fname, font_size=9.5, bold=True, color=C_AMBER)
    add_textbox(slide, 6.85, y + 0.26, 6.1, 0.22, fdesc, font_size=9, color=C_WHITE)
    y += 0.57

add_textbox(slide, 6.7, 5.05, 6.3, 0.28, "ルーティング決定 — analysis_target の値", font_size=10, bold=True, color=C_ACCENT)
table_row(slide, 6.7, 5.38, [2.0, 2.1, 2.2],
          ["analysis_target", "エージェント3・4の取得先", "使用条件"],
          [C_DARK_BLUE, C_DARK_BLUE, C_DARK_BLUE], heights=0.3, bold_flags=[True, True, True])
routing_rows = [
    ("Agent",     "Langfuseのみ",           "エラーがAIエージェントロジック/LLM呼び出しに起因"),
    ("InfraLogs", "Prometheusのみ",          "エラーがインフラ/リソースに起因"),
    ("Unknown",   "Langfuse + Prometheus両方","エラー起因不明 — 両方でクロスチェック"),
]
y = 5.68
bgs = [C_BOX_BG, C_BOX_BG2, C_BOX_BG]
for i, (target, fetch, when) in enumerate(routing_rows):
    table_row(slide, 6.7, y, [2.0, 2.1, 2.2], [target, fetch, when],
              [bgs[i], bgs[i], bgs[i]], heights=0.3)
    y += 0.3


# ══════════════════════════════════════════════════════════════════════════
# スライド 9 — エージェント3: エラー分析
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "エージェント3 — エラー分析エージェント",
             "目的: 詳細なエラー抽出 — 各エラーにユニークID・カテゴリ・重大度・証拠を付与")

add_rect(slide, 0.3, 1.0, 12.7, 0.55, C_BOX_BG)
add_rect(slide, 0.3, 1.0, 0.08, 0.55, C_ACCENT)
add_textbox(slide, 0.5, 1.06, 12.3, 0.42,
            "厳格なスコープ: エラーの識別と分類のみ。根本原因・推奨・推測は行わない。",
            font_size=10, bold=True, color=C_WHITE)

add_textbox(slide, 0.3, 1.7, 5.9, 0.28, "ターゲットデータ取得", font_size=10, bold=True, color=C_ACCENT)

table_row(slide, 0.3, 2.02, [2.0, 1.95, 1.95],
          ["analysis_target", "Langfuse", "Prometheus"],
          [C_DARK_BLUE]*3, heights=0.3, bold_flags=[True, True, True])
fetch_rows = [
    ("Agent",     "✓  取得（trace_id必要）", "✗  取得しない"),
    ("InfraLogs", "✗  取得しない",            "✓  取得"),
    ("Unknown",   "✓  取得（trace_id必要）", "✓  取得"),
]
y = 2.32
for i, (target, lf, pr) in enumerate(fetch_rows):
    bg = C_BOX_BG if i % 2 == 0 else C_BOX_BG2
    table_row(slide, 0.3, y, [2.0, 1.95, 1.95], [target, lf, pr], [bg]*3, heights=0.3)
    y += 0.3

add_textbox(slide, 0.3, 3.05, 5.9, 0.28, "処理ステップ", font_size=10, bold=True, color=C_ACCENT)
steps = [
    "CorrelationResultからanalysis_targetを読み取る",
    "対象データソースのみを取得",
    "全ログをマージして時系列ソート",
    "プロンプト構築: 相関コンテキスト + 正規化データ + ログ",
    "ErrorAnalysisResultスキーマでGPT-4o呼び出し",
    "Pydanticバリデーション — エラー証拠の欠落を拒否",
    "ErrorAnalysisResult + rca_target（パススルー）を返却",
]
y = 3.35
for i, s in enumerate(steps):
    add_rect(slide, 0.3, y, 0.32, 0.28, C_ACCENT)
    add_textbox(slide, 0.3, y + 0.02, 0.32, 0.24, str(i+1),
                font_size=9.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.67, y + 0.02, 5.4, 0.26, s, font_size=9.5, color=C_WHITE)
    y += 0.33

add_textbox(slide, 6.7, 1.7, 6.3, 0.28, "出力: ErrorAnalysisResult", font_size=10, bold=True, color=C_ACCENT)

add_rect(slide, 6.7, 2.02, 6.3, 0.28, C_DARK_BLUE)
add_textbox(slide, 6.85, 2.05, 6.1, 0.22, "トップレベルフィールド", font_size=9, bold=True, color=C_LIGHT_GREY)
top_fields = [
    ("analysis_summary",       "str（最大500文字）"),
    ("analysis_target",        "Agent | InfraLogs | Unknown"),
    ("errors",                 "list[ErrorDetail]  — 最低1件必須"),
    ("error_patterns",         "list[ErrorPattern]  — 2件以上の発生が必要"),
    ("error_impacts",          "list[ErrorImpact]  — コンポーネント毎の影響評価"),
    ("error_propagation_path", "list[str]  — 時系列順の伝播経路"),
    ("confidence",             "float  0.0 – 1.0"),
]
y = 2.3
for fname, ftype in top_fields:
    add_rect(slide, 6.7, y, 6.3, 0.3, C_BOX_BG)
    add_textbox(slide, 6.85, y + 0.04, 2.5, 0.22, fname, font_size=9.5, bold=True, color=C_AMBER)
    add_textbox(slide, 9.4, y + 0.04, 3.5, 0.22, ftype, font_size=9.5, color=C_WHITE)
    y += 0.32

add_textbox(slide, 6.7, 4.55, 6.3, 0.28, "ErrorDetailスキーマ（エラー単位）", font_size=10, bold=True, color=C_ACCENT)
add_rect(slide, 6.7, 4.85, 6.3, 1.85, C_BOX_BG2)
detail_fields = [
    ("error_id",      '"ERR-001"、"ERR-002" ...  ← ユニークな連番ID'),
    ("category",      "llm_failure | timeout | dns_failure | configuration_error | ..."),
    ("severity",      "critical | high | medium | low | info"),
    ("component",     "ログから抽出した正確なサービス/コンポーネント名"),
    ("error_message", "ログから抽出した正確なエラーテキスト"),
    ("evidence",      "このエラーの存在を証明する生ログ行"),
    ("source",        "langfuse | prometheus"),
]
y = 4.9
for fname, fdesc in detail_fields:
    add_textbox(slide, 6.85, y, 1.6, 0.24, fname, font_size=9, bold=True, color=C_AMBER)
    add_textbox(slide, 8.5, y, 4.4, 0.24, fdesc, font_size=9, color=C_WHITE)
    y += 0.26

add_rect(slide, 6.7, 6.75, 6.3, 0.45, C_GREEN)
add_textbox(slide, 6.85, 6.8, 6.1, 0.35,
            "error_idクロス参照: RCAは特定エラーをIDで参照。推奨エージェントも各ソリューションで特定エラーをIDで参照。",
            font_size=9, bold=True, color=C_NAVY)


# ══════════════════════════════════════════════════════════════════════════
# スライド 10 — エージェント4: 根本原因分析 概要
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "エージェント4 — 根本原因分析エージェント",
             "目的: エラーが発生した「なぜ」を決定的に特定 — 根本原因・因果連鎖・5つのなぜ分析")

add_rect(slide, 0.3, 1.0, 12.7, 0.55, C_BOX_BG)
add_rect(slide, 0.3, 1.0, 0.08, 0.55, C_RED)
add_textbox(slide, 0.5, 1.06, 12.3, 0.42,
            "最も分析負荷の高いエージェント。厳格なスコープ: 根本原因の特定のみ。推奨・修復手順は行わない。",
            font_size=10, bold=True, color=C_WHITE)

add_textbox(slide, 0.3, 1.7, 6.0, 0.28, "処理ステップ", font_size=10, bold=True, color=C_ACCENT)
steps = [
    ("1", "rca_target読み取り",      "エラー分析からパススルー（= analysis_target）"),
    ("2", "データ取得ルーティング",  "Agent→Langfuseのみ  |  InfraLogs→Prometheusのみ  |  Unknown→両方"),
    ("3", "新鮮なログの再取得",      "同一ソースから独立取得 — 因果関係分析のための完全なコンテキスト"),
    ("4", "リッチユーザーメッセージ構築","全ErrorDetail + 新鮮なログ + インシデントコンテキスト"),
    ("5", "システムプロンプト構築",  "エラーカテゴリ・伝播経路・5つのなぜルール・RCAResultスキーマ"),
    ("6", "GPT-4o呼び出し",          "temperature=0.0  ·  5つのなぜを強制した完全RCAResult"),
    ("7", "Pydanticバリデーション",  "RCAResult.model_validate — five_why_analysis欠損・不完全は拒否"),
]
y = 2.02
for num, title, desc in steps:
    add_rect(slide, 0.3, y, 0.35, 0.56, C_ACCENT)
    add_textbox(slide, 0.3, y + 0.1, 0.35, 0.35, num,
                font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.65, y, 5.55, 0.56, C_BOX_BG)
    add_textbox(slide, 0.8, y + 0.04, 5.3, 0.24, title, font_size=10, bold=True, color=C_AMBER)
    add_textbox(slide, 0.8, y + 0.3, 5.3, 0.24, desc, font_size=9.5, color=C_LIGHT_GREY)
    y += 0.63

add_textbox(slide, 6.7, 1.7, 6.3, 0.28, "出力: RCAResult", font_size=10, bold=True, color=C_ACCENT)

rca_fields = [
    ("rca_summary",          "str  最大800文字 — エグゼクティブサマリー"),
    ("root_cause",           "RootCause  （カテゴリ・コンポーネント・説明・証拠・error_ids・信頼度）"),
    ("causal_chain",         "list[CausalLink]  最低1件  —  link_type+証拠を持つsource→target"),
    ("contributing_factors", "list[ContributingFactor]  —  障害を悪化させた増幅要因"),
    ("failure_timeline",     "list[FailureTimeline]  —  is_root_causeフラグ付き時系列イベント"),
    ("blast_radius",         "list[str]  —  根本原因の影響を受けた全コンポーネント"),
    ("five_why_analysis",    "FiveWhyAnalysis  —  正確に5つのWhyStep + fundamental_root_cause  ← 新規"),
    ("confidence",           "float  0.0 – 1.0  —  証拠の強度に連動"),
]
y = 2.02
for fname, fdesc in rca_fields:
    is_new = "新規" in fdesc
    bg = C_RED if is_new else C_BOX_BG
    add_rect(slide, 6.7, y, 6.3, 0.56, bg)
    add_textbox(slide, 6.85, y + 0.04, 2.3, 0.24, fname, font_size=9.5, bold=True, color=C_AMBER)
    add_textbox(slide, 6.85, y + 0.3, 6.1, 0.22, fdesc, font_size=9, color=C_WHITE)
    y += 0.6

add_rect(slide, 6.7, 6.82, 6.3, 0.38, C_AMBER)
add_textbox(slide, 6.85, 6.87, 6.1, 0.28,
            "CausalLinkタイプ:  direct_cause  |  indirect_cause  |  trigger  |  amplifier",
            font_size=9.5, bold=True, color=C_NAVY)


# ══════════════════════════════════════════════════════════════════════════
# スライド 11 — エージェント4: 5つのなぜ分析
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "エージェント4 — 5つのなぜ分析",
             "反復的手法: 各回答が次の「なぜ？」の対象になる — 症状から根本原因へ掘り下げる")

add_textbox(slide, 0.3, 1.0, 12.7, 0.28, "WhyStepスキーマ（5件必須）", font_size=10, bold=True, color=C_ACCENT)
schema_fields = [
    ("step",      "int 1–5",  "シーケンス内の位置"),
    ("question",  "str",      '"なぜ[前回の回答]が発生したのか？"'),
    ("answer",    "str",      "このレベルでの原因の説明"),
    ("evidence",  "str",      "回答を裏付ける具体的なログ行またはメトリクス"),
    ("component", "str",      "関与するサービスまたはシステム要素"),
]
table_row(slide, 0.3, 1.32, [1.4, 1.4, 9.2],
          ["フィールド", "型", "説明"],
          [C_DARK_BLUE]*3, heights=0.28, bold_flags=[True]*3)
y = 1.6
for fname, ftype, fdesc in schema_fields:
    bg = C_BOX_BG if y % 0.6 > 0.3 else C_BOX_BG2
    table_row(slide, 0.3, y, [1.4, 1.4, 9.2], [fname, ftype, fdesc], [bg]*3, heights=0.28)
    y += 0.28

add_textbox(slide, 0.3, 2.9, 12.7, 0.28,
            "例  —  'medical-rag'サービスでのLLMアクセス無効化", font_size=10, bold=True, color=C_ACCENT)

example_whys = [
    ("問題",    "medical-ragがリクエスト処理に失敗 — LLMアクセスが無効化されている",
     "LLMアクセスが無効（デモエラーモード）のログエントリ"),
    ("なぜ1",   "サービスのLLMアクセスが無効化されている",
     "ログ: LLM access is disabled (demo error mode). Click Enable LLM Access..."),
    ("なぜ2",   "サービスがデモエラーモードで動作するよう設定されている",
     "エラーメッセージが原因として明示的に'demo error mode'を参照"),
    ("なぜ3",   "デモエラーモードの設定フラグが切り替えられた（UIまたはデプロイ経由）",
     "エラーがUIの'Enable LLM Access'クリックを指示 — UIトグルがこれを制御"),
    ("なぜ4",   "リクエスト受付前にLLMアクセス状態を確認するプリフライトバリデーションがない",
     "サービスはリクエストを受け付け、LLM呼び出し段階でのみ失敗 — 早期拒否ログなし"),
    ("なぜ5",   "対応するレディネスゲートなしでデモエラーモードがテスト用に追加された",
     "デモモードを拒否するスタートアップ/レディネスログの不在から推定（可視性の限界）"),
]

colors_ex = [C_DARK_BLUE, C_BOX_BG, C_BOX_BG, C_BOX_BG, C_BOX_BG, C_RED]
label_colors = [C_LIGHT_GREY, C_WHITE, C_WHITE, C_WHITE, C_WHITE, C_AMBER]

y = 3.22
for i, (label, answer, evidence) in enumerate(example_whys):
    add_rect(slide, 0.3, y, 1.2, 0.55, colors_ex[i])
    add_textbox(slide, 0.3, y + 0.12, 1.2, 0.3, label,
                font_size=10, bold=True, color=label_colors[i], align=PP_ALIGN.CENTER)
    add_rect(slide, 1.5, y, 11.5, 0.55, C_BOX_BG if i % 2 == 0 else C_BOX_BG2)
    add_textbox(slide, 1.65, y + 0.02, 11.2, 0.24, answer, font_size=9.5, bold=(i==5), color=C_WHITE)
    add_textbox(slide, 1.65, y + 0.3, 11.2, 0.22, f"証拠: {evidence}", font_size=8.5, color=C_LIGHT_GREY, italic=True)
    y += 0.6

add_rect(slide, 0.3, 6.88, 12.7, 0.38, C_GREEN)
add_textbox(slide, 0.5, 6.93, 12.3, 0.28,
            "根本的な根本原因:  デモエラーモード機能にレディネスゲートが欠如 — 設定トグルに強制メカニズムが対応付けられていなかった",
            font_size=9.5, bold=True, color=C_NAVY)

add_rect(slide, 0.3, y + 0.05, 12.7, 0.0, C_BOX_BG)


# ══════════════════════════════════════════════════════════════════════════
# スライド 12 — エージェント5: 推奨
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "エージェント5 — 推奨エージェント",
             "目的: エラー分析と根本原因分析を統合し、1〜4件のランク付き実行可能ソリューションを生成")

add_rect(slide, 0.3, 1.0, 12.7, 0.52, C_GREEN)
add_textbox(slide, 0.5, 1.07, 12.3, 0.38,
            "外部データを一切取得しない唯一のエージェント。純粋な統合レイヤー — Langfuse・Prometheus・I/O待機なし。パイプラインで最速のステップ。",
            font_size=10, bold=True, color=C_NAVY)

add_textbox(slide, 0.3, 1.68, 6.0, 0.28, "処理ステップ", font_size=10, bold=True, color=C_ACCENT)
steps = [
    ("1", "リッチユーザーメッセージ構築",  "RCAサマリー + 因果連鎖 + 寄与要因 + 影響範囲 + 全エラー + パターン + 影響"),
    ("2", "システムプロンプト構築",         "RCA(8件)+エラー分析(5件)から13フィールドを注入 + ランク付けルール + ソリューションスキーマ"),
    ("3", "GPT-4o呼び出し",                "temperature=0.0  ·  1〜4件のみ — パディング防止ルールをプロンプトで強制"),
    ("4", "Pydanticバリデーション",         "model_validatorがランクの1〜N連番を確認 — ギャップや重複は拒否"),
]
y = 2.0
for num, title, desc in steps:
    add_rect(slide, 0.3, y, 0.35, 0.88, C_ACCENT)
    add_textbox(slide, 0.3, y + 0.22, 0.35, 0.4, num,
                font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.65, y, 5.55, 0.88, C_BOX_BG)
    add_textbox(slide, 0.8, y + 0.08, 5.3, 0.28, title, font_size=10, bold=True, color=C_AMBER)
    add_textbox(slide, 0.8, y + 0.42, 5.3, 0.4, desc, font_size=9.5, color=C_LIGHT_GREY)
    y += 0.98

add_textbox(slide, 0.3, 5.65, 6.0, 0.28, "ランク付けルール", font_size=10, bold=True, color=C_ACCENT)
ranking = [
    ("ランク1", C_RED,    "根本原因を直接修正 — addresses_root_cause=True"),
    ("ランク2", C_AMBER,  "最も重大な二次的懸念または伝播防止"),
    ("ランク3", C_ACCENT, "寄与要因または回復力の改善"),
    ("ランク4", C_GREEN,  "将来の再発防止策（本当に有益な場合のみ）"),
]
y = 5.95
for rank, color, rule in ranking:
    add_rect(slide, 0.3, y, 1.0, 0.3, color)
    add_textbox(slide, 0.3, y + 0.04, 1.0, 0.22, rank,
                font_size=9.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.35, y + 0.04, 4.85, 0.24, rule, font_size=9.5, color=C_WHITE)
    y += 0.35

add_textbox(slide, 6.7, 1.68, 6.3, 0.28, "出力: Solutionスキーマ", font_size=10, bold=True, color=C_ACCENT)
solution_fields = [
    ("rank",                "int 1–4  — ユニーク・連番（Pydantic検証）"),
    ("title",               "str  最大120文字  — 短い実行可能タイトル"),
    ("description",         "str  — 詳細アクション + 根本原因への対処理由"),
    ("category",            "config_change | code_fix | infrastructure | scaling | retry_logic | fallback | monitoring | ..."),
    ("effort",              "quick_fix | low | medium | high"),
    ("addresses_root_cause","bool  — 根本原因を直接修正するソリューションのみTrue"),
    ("affected_components", "list[str]  — このソリューションが対象とするコンポーネント"),
    ("expected_outcome",    "str  — 実装後に期待される改善"),
    ("error_ids",           "list[str]  — ErrorAnalysisのエラーIDへの参照（ERR-001等）"),
]
y = 2.0
for fname, fdesc in solution_fields:
    add_rect(slide, 6.7, y, 6.3, 0.52, C_BOX_BG if y % 1.0 > 0.5 else C_BOX_BG2)
    add_textbox(slide, 6.85, y + 0.03, 2.4, 0.22, fname, font_size=9.5, bold=True, color=C_AMBER)
    add_textbox(slide, 6.85, y + 0.28, 6.1, 0.2, fdesc, font_size=9, color=C_WHITE)
    y += 0.56


# ══════════════════════════════════════════════════════════════════════════
# スライド 13 — パイプラインデータフロー
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "エンドツーエンド パイプラインデータフロー",
             "各エージェントが受け取るデータ、外部から取得するデータ、次のエージェントへ渡すデータ")

flow = [
    ("正規化",         "タイムスタンプ\ntrace_id?\nagent_name",
     "Langfuse（trace_idあり）\nPrometheus（trace_idなし）",
     "NormalizedIncident\nerror_type、signals、\nentities、confidence"),
    ("相関分析",       "NormalizedIncident\ntrace_id?\nagent_name",
     "Prometheus（常時）\nLangfuse（trace_idあり）",
     "CorrelationResult\n因果連鎖、タイムライン、\nanalysis_target  ←重要"),
    ("エラー分析",     "CorrelationResult\nNormalizedIncident\ntrace_id? + agent_name",
     "Langfuse（Agent/Unk時）\nPrometheus（Infra/Unk時）",
     "ErrorAnalysisResult\nID付きエラー、\nパターン、rca_target"),
    ("根本原因分析",   "ErrorAnalysisResult\nNormalizedIncident\nrca_target + trace_id?",
     "Langfuse（Agent/Unk時）\nPrometheus（Infra/Unk時）",
     "RCAResult\nroot_cause、causal_chain\nfive_why_analysis"),
    ("推奨",           "ErrorAnalysisResult\nRCAResult\nagent_name",
     "なし — 統合のみ",
     "RecommendationResult\n1〜4件のランク付きソリューション\n工数 + カテゴリ"),
]

table_row(slide, 0.3, 1.0, [2.1, 3.2, 3.4, 4.0],
          ["エージェント", "前エージェントから受け取る", "外部取得", "次エージェントへ渡す"],
          [C_DARK_BLUE]*4, heights=0.35, bold_flags=[True]*4)

y = 1.35
row_colors = [C_BOX_BG, C_BOX_BG2, C_BOX_BG, C_BOX_BG2, C_BOX_BG]
for i, (agent, receives, fetches, passes) in enumerate(flow):
    h = 1.05
    fetch_color = C_BOX_BG2 if "なし" in fetches else row_colors[i]
    fc_text = C_GREEN if "なし" in fetches else C_WHITE

    add_rect(slide, 0.3, y, 2.1, h, C_ACCENT)
    add_textbox(slide, 0.3, y + 0.3, 2.1, 0.4, agent,
                font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, 2.4, y, 3.2, h, row_colors[i])
    add_textbox(slide, 2.5, y + 0.08, 3.0, h - 0.15, receives, font_size=9, color=C_WHITE)

    add_rect(slide, 5.6, y, 3.4, h, fetch_color)
    add_textbox(slide, 5.7, y + 0.08, 3.2, h - 0.15, fetches, font_size=9, color=fc_text)

    add_rect(slide, 9.0, y, 4.1, h, row_colors[i])
    add_textbox(slide, 9.1, y + 0.08, 3.9, h - 0.15, passes, font_size=9, color=C_WHITE)

    y += h + 0.04

add_textbox(slide, 0.3, 6.92, 12.7, 0.28,
            "外部データは必要な各エージェントが独立して取得 — 共有ログキャッシュなし。再取得により各エージェントが特定の分析タスクに対して完全かつフィルタリングされていない証拠を得られる。",
            font_size=9, color=C_LIGHT_GREY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# スライド 14 — 効率性・設計方針
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "効率性と主要な設計方針")

decisions = [
    (C_GREEN,  "NO_ERROR短絡処理",
     "正規化の前LLMスキャン: エラー信号未検出 → 即時返却、全5エージェントをスキップ。最も頻繁な本番ケース（正常トレース）。実行ごとに約$0.10〜$0.30を節約。"),
    (C_ACCENT, "スキーマ駆動LLM契約",
     "全Pydanticモデルがmodel_json_schema()でJSONスキーマを自動生成し、システムプロンプトに注入。新しいフィールドの追加がLLM契約を自動更新 — 手動JSONメンテナンス不要。"),
    (C_AMBER,  "エラーIDクロス参照",
     "エラー分析が各個別エラーにERR-001、ERR-002...を割り当て。RCAは因果連鎖とroot_cause.error_idsでこれらのIDを参照。推奨は各solution.error_idsで参照。完全なエラーオブジェクトを再埋め込みせず軽量にリンク。"),
    (C_RED,    "ルーティング決定の前方伝播",
     "相関分析がanalysis_targetを一度設定。エラー分析がrca_targetとしてパススルー。RCAがrca_targetを読み取る。各下流エージェントが状況を再分析することなくルーティング決定を継承。"),
    (C_GREEN,  "グレースフルデグラデーション",
     "全ての外部ソース障害（Langfuse・Prometheus）がWARNプレースホルダーエントリを生成。LLMにデータ欠損を明示し、信頼度を低く設定するよう指示。エージェントはエラーを返さず、パイプラインは常に完了する。"),
    (C_ACCENT, "5つのなぜスキーマ強制",
     "FiveWhyAnalysis.whysがPydanticでmin_length=5、max_length=5。LLMは5ステップ未満も超過も返却不可。プロンプトが正確な連鎖方法論を定義。各ステップで証拠が必須。"),
]

y = 1.05
for i, (color, title, desc) in enumerate(decisions):
    col = i % 2
    row = i // 2
    lx = 0.3 + col * 6.55
    ly = y + row * 1.85
    add_rect(slide, lx, ly, 6.3, 1.72, C_BOX_BG)
    add_rect(slide, lx, ly, 0.08, 1.72, color)
    add_textbox(slide, lx + 0.22, ly + 0.1, 5.95, 0.3, title, font_size=11, bold=True, color=color)
    add_textbox(slide, lx + 0.22, ly + 0.45, 5.95, 1.15, desc, font_size=9.5, color=C_WHITE)


# ══════════════════════════════════════════════════════════════════════════
# スライド 15 — まとめ クイックリファレンス
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "まとめ — クイックリファレンス", "5つのエージェントを一覧で確認")

cols = [2.2, 2.0, 2.5, 2.5, 3.8]
table_row(slide, 0.15, 1.0, cols,
          ["エージェント", "データソース", "主要出力", "重要フィールド", "スコープ境界"],
          [C_DARK_BLUE]*5, heights=0.35, bold_flags=[True]*5)

summary_rows = [
    ("1  正規化",       "Langfuse または Prometheus\n（相互排他 — 両方同時は使用しない）",
     "NormalizedIncident", "error_type",      "エラー分類のみ — 因果関係なし"),
    ("2  相関分析",     "Prometheus（常時）\n+ Langfuse（trace_idあり）",
     "CorrelationResult",  "analysis_target", "因果グラフのみ — エラー詳細なし"),
    ("3  エラー分析",   "Langfuse または Prometheus\n（analysis_targetに基づく）",
     "ErrorAnalysisResult","error_ids (ERR-N)","エラー抽出のみ — 根本原因なし"),
    ("4  根本原因分析", "Langfuse または Prometheus\n（rca_targetに基づく）",
     "RCAResult",          "five_why_analysis","根本原因のみ — 推奨なし"),
    ("5  推奨",         "なし\n（統合のみ）",
     "RecommendationResult","rank + error_ids","ソリューションのみ — 診断なし"),
]
y = 1.35
row_bgs = [C_BOX_BG, C_BOX_BG2, C_BOX_BG, C_BOX_BG2, C_BOX_BG]
for i, (agent, sources, output, key_field, boundary) in enumerate(summary_rows):
    h = 0.88
    table_row(slide, 0.15, y, cols,
              [agent, sources, output, key_field, boundary],
              [row_bgs[i]]*5, heights=h)
    y += h + 0.02

stats = [
    ("5", "パイプライン内\nエージェント数"),
    ("2", "外部データソース\n（Langfuse + Prometheus）"),
    ("6", "Prometheusフェッチ毎の\nPromQLクエリ数"),
    ("0.0", "LLM temperature\n（決定論的出力）"),
    ("5", "5つのなぜ分析の\nWhyStepレスポンス数"),
    ("1–4", "推奨エージェントの\nランク付きソリューション数"),
]
x = 0.15
for stat, label in stats:
    add_rect(slide, x, 5.95, 2.15, 1.3, C_BOX_BG)
    add_textbox(slide, x, 6.02, 2.15, 0.65, stat,
                font_size=30, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, 6.67, 2.15, 0.5, label,
                font_size=8.5, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER)
    x += 2.2

add_textbox(slide, 0.3, 7.18, 12.7, 0.28,
            "調査パイプライン  ·  マルチエージェント可観測性システム  ·  Prodapt AIチーム  ·  2026",
            font_size=9, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)


# ── 保存 ───────────────────────────────────────────────────────────────────
output_path = r"C:\Users\vyanktesh.l\Documents\Invastigate_flow_with_Poller\doc_for_reference\Agent_Pipeline_Documentation_JA.pptx"
prs.save(output_path)
print(f"保存完了: {output_path}")
print(f"スライド数: {len(prs.slides)}")
