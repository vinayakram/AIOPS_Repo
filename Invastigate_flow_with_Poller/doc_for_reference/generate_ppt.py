"""
Generate Investigation Pipeline Agent Documentation PPT
Run: python doc_for_reference/generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ─────────────────────────────────────────────────────────
C_NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # slide background
C_DARK_BLUE  = RGBColor(0x1A, 0x35, 0x5E)   # section header bar
C_ACCENT     = RGBColor(0x2E, 0x86, 0xAB)   # accent blue (titles)
C_GREEN      = RGBColor(0x27, 0xAE, 0x60)   # success / data source green
C_AMBER      = RGBColor(0xF3, 0x9C, 0x12)   # highlight / warning amber
C_RED        = RGBColor(0xC0, 0x39, 0x2B)   # fail / critical
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)   # body text
C_LIGHT_GREY = RGBColor(0xCC, 0xD6, 0xE0)   # sub-text / dividers
C_BOX_BG     = RGBColor(0x16, 0x27, 0x3E)   # content box background
C_BOX_BG2    = RGBColor(0x1E, 0x34, 0x52)   # alternate box background
C_STEP_BG    = RGBColor(0x0A, 0x42, 0x6B)   # step bubble background


# ── Helpers ────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=12, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def slide_header(slide, title, subtitle=None, accent_bar=True):
    """Adds a consistent header bar + title to a content slide."""
    if accent_bar:
        bar = add_rect(slide, 0, 0, 13.33, 0.75, C_DARK_BLUE)
        add_textbox(slide, 0.3, 0.1, 12.5, 0.55, title,
                    font_size=20, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, 0.3, 0.78, 12.5, 0.35, subtitle,
                    font_size=11, color=C_LIGHT_GREY, italic=True)


def bullet_block(slide, left, top, width, height, title, bullets,
                 title_color=C_ACCENT, bg=C_BOX_BG, bullet_color=C_WHITE,
                 title_size=13, bullet_size=10.5):
    """Creates a labelled bullet box."""
    add_rect(slide, left, top, width, height, bg)
    add_textbox(slide, left + 0.15, top + 0.05, width - 0.2, 0.3,
                title, font_size=title_size, bold=True, color=title_color)
    y = top + 0.35
    for b in bullets:
        add_textbox(slide, left + 0.25, y, width - 0.35, 0.28,
                    f"• {b}", font_size=bullet_size, color=bullet_color)
        y += 0.28
    return y


def step_box(slide, left, top, width, num, title, desc, num_color=C_AMBER):
    """Renders a numbered step box."""
    add_rounded_rect(slide, left, top, width, 0.72, C_STEP_BG, C_ACCENT)
    add_textbox(slide, left + 0.08, top + 0.06, 0.35, 0.55,
                str(num), font_size=22, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.48, top + 0.06, width - 0.58, 0.28,
                title, font_size=11, bold=True, color=C_WHITE)
    add_textbox(slide, left + 0.48, top + 0.33, width - 0.58, 0.35,
                desc, font_size=9.5, color=C_LIGHT_GREY)


def table_row(slide, left, top, cols, values, colors, heights=0.32, bold_flags=None):
    """Renders a simple table row using rectangles and text."""
    x = left
    for i, (col_w, val) in enumerate(zip(cols, values)):
        add_rect(slide, x, top, col_w, heights, colors[i])
        bf = bold_flags[i] if bold_flags else False
        add_textbox(slide, x + 0.08, top + 0.04, col_w - 0.1, heights - 0.05,
                    str(val), font_size=9.5, bold=bf, color=C_WHITE)
        x += col_w


def arrow_right(slide, left, top, width=0.4):
    """Adds a right-pointing arrow connector."""
    add_textbox(slide, left, top, width, 0.25, "→",
                font_size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def arrow_down(slide, left, top, height=0.25):
    add_textbox(slide, left, top, 0.35, height, "▼",
                font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ── Presentation setup ─────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)

# Large accent stripe
add_rect(slide, 0, 2.8, 13.33, 0.08, C_ACCENT)

add_textbox(slide, 1.0, 0.9, 11.33, 1.2,
            "Investigation Pipeline",
            font_size=42, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 1.0, 2.05, 11.33, 0.7,
            "Multi-Agent Observability System — Technical Documentation",
            font_size=20, color=C_ACCENT, align=PP_ALIGN.CENTER)

agents = [
    ("1", "Normalization"), ("2", "Correlation"), ("3", "Analysis"),
    ("4", "RCA"), ("5", "Recommendation"),
]
x = 1.0
for num, name in agents:
    add_rounded_rect(slide, x, 3.2, 2.1, 0.9, C_DARK_BLUE, C_ACCENT)
    add_textbox(slide, x, 3.22, 2.1, 0.38, num,
                font_size=16, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, 3.58, 2.1, 0.42, name,
                font_size=11, color=C_WHITE, align=PP_ALIGN.CENTER)
    if num != "5":
        add_textbox(slide, x + 2.1, 3.5, 0.2, 0.3, "→",
                    font_size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    x += 2.27

add_textbox(slide, 1.0, 4.35, 11.33, 0.4,
            "Langfuse  ·  Prometheus  ·  GPT-4o  ·  Pydantic  ·  FastAPI",
            font_size=13, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 6.9, 12.33, 0.35,
            "Prodapt  ·  AI Observability Team  ·  2026",
            font_size=10, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "Agenda")

items = [
    ("01", "System Overview",          "Architecture, key principles, tech stack"),
    ("02", "Data Sources",             "Langfuse (AI traces) and Prometheus (infra metrics)"),
    ("03", "Agent 1 — Normalization",  "Raw log ingestion, error classification, NO_ERROR short-circuit"),
    ("04", "Agent 2 — Correlation",    "Cross-system causal graph, analysis_target routing"),
    ("05", "Agent 3 — Error Analysis", "Deep-dive error extraction, pattern detection, error IDs"),
    ("06", "Agent 4 — RCA",            "Root cause determination, causal chain, Five Whys analysis"),
    ("07", "Agent 5 — Recommendation", "Ranked actionable solutions, effort estimation"),
    ("08", "Pipeline Data Flow",       "End-to-end data passing between agents"),
    ("09", "Efficiency & Design",      "Schema contracts, short-circuits, graceful degradation"),
]

y = 1.05
for num, title, desc in items:
    add_rect(slide, 0.4, y, 0.5, 0.52, C_ACCENT)
    add_textbox(slide, 0.4, y + 0.08, 0.5, 0.35, num,
                font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.9, y, 11.9, 0.52, C_BOX_BG)
    add_textbox(slide, 1.05, y + 0.04, 3.8, 0.26, title,
                font_size=12, bold=True, color=C_WHITE)
    add_textbox(slide, 1.05, y + 0.28, 11.0, 0.22, desc,
                font_size=9.5, color=C_LIGHT_GREY)
    y += 0.6


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — System Overview
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "System Overview", "Sequential 5-agent pipeline — each agent has one clearly scoped responsibility")

# Pipeline flow boxes
agents_flow = [
    ("1", "Normalization",   "Classify error type\nExtract signals", C_BOX_BG),
    ("2", "Correlation",     "Build causal graph\nRoute analysis", C_BOX_BG),
    ("3", "Error Analysis",  "Extract & ID errors\nDetect patterns", C_BOX_BG),
    ("4", "RCA",             "Root cause +\nFive Whys analysis", C_BOX_BG),
    ("5", "Recommendation",  "Ranked actionable\nsolutions (1–4)", C_BOX_BG),
]

x = 0.3
for num, name, desc, bg in agents_flow:
    add_rounded_rect(slide, x, 1.1, 2.35, 1.6, bg, C_ACCENT)
    add_rect(slide, x, 1.1, 2.35, 0.38, C_ACCENT)
    add_textbox(slide, x, 1.12, 2.35, 0.35,
                f"Step {num}  •  {name}", font_size=11, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 1.55, 2.15, 0.9, desc,
                font_size=9.5, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER)
    if num != "5":
        add_textbox(slide, x + 2.35, 1.72, 0.25, 0.3, "→",
                    font_size=18, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    x += 2.6

# Key principles
principles = [
    ("Deterministic", "temperature=0.0 on all LLM calls — reproducible outputs"),
    ("Schema-Driven", "Pydantic models auto-generate the JSON schema injected into every prompt"),
    ("Graceful", "Every external failure produces a WARN placeholder — no agent crashes"),
    ("Isolated", "Each agent sees only the previous agent's output — no shared state"),
]
y = 3.0
add_textbox(slide, 0.3, 2.9, 12.7, 0.3, "KEY PRINCIPLES",
            font_size=10, bold=True, color=C_ACCENT)
for i, (title, desc) in enumerate(principles):
    col = i % 2
    row = i // 2
    lx = 0.3 + col * 6.55
    ly = y + 0.05 + row * 0.95
    add_rect(slide, lx, ly, 6.3, 0.82, C_BOX_BG2)
    add_rect(slide, lx, ly, 0.08, 0.82, C_GREEN)
    add_textbox(slide, lx + 0.2, ly + 0.06, 6.0, 0.28,
                title, font_size=11, bold=True, color=C_GREEN)
    add_textbox(slide, lx + 0.2, ly + 0.36, 6.0, 0.38,
                desc, font_size=9.5, color=C_LIGHT_GREY)

add_textbox(slide, 0.3, 6.85, 12.5, 0.28,
            "Tech Stack:  GPT-4o  ·  FastAPI  ·  Pydantic v2  ·  AsyncOpenAI  ·  httpx  ·  SQLite",
            font_size=9, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Data Sources: Langfuse
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "Data Source 1 — Langfuse (AI Agent Traces)")

# Left: what + how
add_rect(slide, 0.3, 1.0, 5.9, 5.8, C_BOX_BG)
add_textbox(slide, 0.5, 1.05, 5.6, 0.3, "WHAT IT IS",
            font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 0.5, 1.35, 5.6, 0.55,
            "LLM observability platform that records every AI agent interaction as traces, spans, and generations",
            font_size=10, color=C_WHITE)

add_textbox(slide, 0.5, 1.95, 5.6, 0.3, "AUTHENTICATION",
            font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 0.5, 2.25, 5.6, 0.3,
            "HTTP Basic Auth  →  public_key : secret_key",
            font_size=10, color=C_LIGHT_GREY)

add_textbox(slide, 0.5, 2.65, 5.6, 0.3, "API CALLS PER FETCH (2 total)",
            font_size=10, bold=True, color=C_ACCENT)

calls = [
    ("GET", "/api/public/traces/{trace_id}", "Top-level trace: name, status, input, output, latency, cost"),
    ("GET", "/api/public/observations?traceId=...&limit=100", "All child spans/generations with per-span status, model, tokens, errors"),
]
y = 2.98
for method, endpoint, desc in calls:
    add_rect(slide, 0.5, y, 0.55, 0.28, C_GREEN)
    add_textbox(slide, 0.5, y + 0.02, 0.55, 0.24, method,
                font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.1, y, 4.7, 0.28, endpoint,
                font_size=9, color=C_AMBER)
    add_textbox(slide, 0.5, y + 0.3, 5.6, 0.28, desc,
                font_size=9, color=C_LIGHT_GREY)
    y += 0.68

add_textbox(slide, 0.5, 4.4, 5.6, 0.3, "ACCESSED BY",
            font_size=10, bold=True, color=C_ACCENT)
agents_using = ["Normalization (if trace_id provided)", "Correlation (always if trace_id available)",
                "Error Analysis (if target = Agent or Unknown)", "RCA (if target = Agent or Unknown)"]
yy = 4.7
for a in agents_using:
    add_textbox(slide, 0.6, yy, 5.4, 0.24, f"• {a}", font_size=9.5, color=C_WHITE)
    yy += 0.26

# Right: error detection logic
add_rect(slide, 6.5, 1.0, 6.55, 5.8, C_BOX_BG2)
add_textbox(slide, 6.7, 1.05, 6.2, 0.3, "ERROR DETECTION LOGIC",
            font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 6.7, 1.35, 6.2, 0.35,
            "Each observation is classified ERROR / WARN / INFO by checking 4 fields in order:",
            font_size=9.5, color=C_WHITE)

detection_steps = [
    ("1", "status field",       "ERROR / FAIL / FAILED → ERROR"),
    ("2", "statusMessage",      "Contains error keyword → ERROR"),
    ("3", "output field",       "JSON with 'error' key → ERROR  |  string match → ERROR"),
    ("4", "input field",        "Error keyword match → WARN"),
    ("5", "Default",            "→ INFO"),
]
y = 1.8
for num, field, rule in detection_steps:
    add_rect(slide, 6.6, y, 0.35, 0.38, C_ACCENT)
    add_textbox(slide, 6.6, y + 0.04, 0.35, 0.3, num,
                font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 7.0, y + 0.02, 1.5, 0.18, field,
                font_size=9.5, bold=True, color=C_AMBER)
    add_textbox(slide, 7.0, y + 0.2, 5.8, 0.2, rule,
                font_size=9, color=C_LIGHT_GREY)
    y += 0.5

add_textbox(slide, 6.7, 4.55, 6.2, 0.3, "OUTPUT FORMAT (per log entry)",
            font_size=10, bold=True, color=C_ACCENT)
fields = ["timestamp", "source  (langfuse)", "service  (observation name)",
          "message  (type + status + error + tokens)", "level  (ERROR / WARN / INFO)",
          "metadata  (ids, model, cost, latency, input/output)"]
y = 4.85
for f in fields:
    add_textbox(slide, 6.8, y, 6.1, 0.24, f"→  {f}", font_size=9.5, color=C_WHITE)
    y += 0.26


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Data Sources: Prometheus
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "Data Source 2 — Prometheus (Infrastructure Metrics)")

add_rect(slide, 0.3, 1.0, 12.7, 1.1, C_BOX_BG)
add_textbox(slide, 0.5, 1.05, 8.0, 0.3, "HOW IT IS QUERIED",
            font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 0.5, 1.35, 12.2, 0.55,
            "GET /api/v1/query_range  ·  PromQL with {agent} templating  ·  ±5 min window around incident timestamp  ·  15-second step resolution  ·  10 sec HTTP timeout",
            font_size=10, color=C_WHITE)

# PromQL table
add_textbox(slide, 0.3, 2.2, 12.7, 0.3, "6 DEFAULT PromQL QUERIES (run for every fetch)",
            font_size=10, bold=True, color=C_ACCENT)

# Header row
cols = [2.2, 5.5, 2.5, 2.3]
table_row(slide, 0.3, 2.55, cols,
          ["Query Name", "PromQL Pattern", "What It Detects", "Severity Rule"],
          [C_DARK_BLUE, C_DARK_BLUE, C_DARK_BLUE, C_DARK_BLUE],
          heights=0.32, bold_flags=[True, True, True, True])

rows = [
    ("error_rate",    "rate(http_requests_total{status=~'5..',job=~'.*{agent}.*'}[window])", "HTTP 5xx error rate",        "ERROR if > 0"),
    ("latency_p99",   "histogram_quantile(0.99, rate(http_request_duration_seconds_bucket...))",  "P99 latency spikes",  "INFO only"),
    ("up_status",     "up{job=~'.*{agent}.*'}",                                              "Service up / down",          "ERROR if == 0"),
    ("memory_usage",  "container_memory_usage_bytes{pod=~'.*{agent}.*'}",                    "Container memory consumption","INFO only"),
    ("restart_count", "kube_pod_container_status_restarts_total{pod=~'.*{agent}.*'}",         "Pod crash-loop detection",   "WARN if > 0"),
    ("dns_failures",  "rate(coredns_dns_responses_total{rcode='SERVFAIL'}[window])",          "DNS SERVFAIL rate",          "ERROR if > 0"),
]
y = 2.87
for i, (name, promql, detects, severity) in enumerate(rows):
    bg = C_BOX_BG if i % 2 == 0 else C_BOX_BG2
    sev_color = C_RED if "ERROR" in severity else C_AMBER if "WARN" in severity else C_BOX_BG
    table_row(slide, 0.3, y, cols,
              [name, promql, detects, severity],
              [bg, bg, bg, sev_color], heights=0.3)
    y += 0.3

# Bottom panels
add_rect(slide, 0.3, 5.6, 4.0, 1.6, C_BOX_BG)
add_textbox(slide, 0.5, 5.65, 3.7, 0.28, "TIME WINDOW", font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 0.5, 5.95, 3.7, 0.9,
            "Incident timestamp ± 5 minutes\nCaptures both the lead-up and aftermath of the failure",
            font_size=10, color=C_WHITE)

add_rect(slide, 4.5, 5.6, 4.2, 1.6, C_BOX_BG)
add_textbox(slide, 4.7, 5.65, 3.9, 0.28, "FAILED QUERY HANDLING", font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 4.7, 5.95, 3.9, 0.9,
            "Each query failure produces a WARN placeholder entry so the LLM knows data was missing — does not abort",
            font_size=10, color=C_WHITE)

add_rect(slide, 8.9, 5.6, 4.1, 1.6, C_BOX_BG)
add_textbox(slide, 9.1, 5.65, 3.8, 0.28, "DATA POINT SELECTION", font_size=10, bold=True, color=C_ACCENT)
add_textbox(slide, 9.1, 5.95, 3.8, 0.9,
            "Only the LAST data point per time series is used — most recent state at the time of the incident",
            font_size=10, color=C_WHITE)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Agent 1: Normalization
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "Agent 1 — Normalization Agent",
             "Purpose: Convert raw, unstructured logs into a single structured NormalizedIncident object")

# Purpose box
add_rect(slide, 0.3, 1.0, 8.6, 0.55, C_BOX_BG)
add_rect(slide, 0.3, 1.0, 0.08, 0.55, C_AMBER)
add_textbox(slide, 0.5, 1.06, 8.3, 0.42,
            "STRICT SCOPE: Classify the error type + extract atomic signals. No causality. No root cause. No correlation.",
            font_size=10, bold=True, color=C_AMBER)

# Steps
add_textbox(slide, 0.3, 1.7, 12.7, 0.28, "INTERNAL PROCESSING STEPS", font_size=10, bold=True, color=C_ACCENT)

steps = [
    ("1", "Route to Data Source",     "trace_id present → Langfuse    |    trace_id absent → Prometheus"),
    ("2", "Fetch Raw Logs",           "Langfuse: 1 trace + up to 100 observations    |    Prometheus: 6 PromQL range queries"),
    ("3", "Pre-LLM Error Scan",       "_has_error_signals(): level check (ERROR/WARN) + whole-word keyword regex scan"),
    ("4", "NO_ERROR Short-Circuit",   "If no signals found → return NO_ERROR immediately, confidence=1.0 — LLM NOT called"),
    ("5", "Build LLM Prompt",         "Format logs as [N] ts=... svc=... level=... msg=... meta={...}  +  inject Pydantic schema"),
    ("6", "Call GPT-4o",              "temperature=0.0  ·  response_format=json_object  ·  schema enforced in prompt"),
    ("7", "Pydantic Validate",        "NormalizedIncident.model_validate(data) — rejects any schema mismatch before returning"),
]
y = 2.05
for num, title, desc in steps:
    step_box(slide, 0.3, y, 12.7, num, title, desc)
    y += 0.78

# Output contract
add_rect(slide, 9.15, 1.0, 4.15, 0.55, C_BOX_BG2)
add_textbox(slide, 9.35, 1.06, 3.9, 0.42, "Routing:  trace_id → Langfuse   |   no trace_id → Prometheus",
            font_size=9.5, color=C_LIGHT_GREY, italic=True)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Agent 1: Output Contract + Key Efficiency
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "Agent 1 — Output Contract & Efficiency Features")

# Output model
add_textbox(slide, 0.3, 1.0, 6.0, 0.3, "OUTPUT: NormalizedIncident", font_size=11, bold=True, color=C_ACCENT)
fields = [
    ("error_type",    "NO_ERROR | AI_AGENT | INFRA | NETWORK | UNKNOWN",   "Priority: INFRA > NETWORK > AI_AGENT when mixed"),
    ("error_summary", "str  (max 300 chars)",                               "Factual 1-2 line description — no inference"),
    ("timestamp",     "ISO-8601 string",                                    "First occurrence extracted from logs"),
    ("confidence",    "float  0.0 – 1.0",                                   "Based on explicitness of log evidence"),
    ("entities",      "{ agent_id, service, trace_id }",                    "Extracted entity references"),
    ("signals",       "list[str]",                                          "e.g. ['LLM_access_disabled', 'timeout']"),
]
y = 1.35
for fname, ftype, fdesc in fields:
    add_rect(slide, 0.3, y, 6.1, 0.62, C_BOX_BG)
    add_textbox(slide, 0.45, y + 0.04, 2.0, 0.24, fname, font_size=10, bold=True, color=C_AMBER)
    add_textbox(slide, 2.5, y + 0.04, 3.8, 0.24, ftype, font_size=10, color=C_WHITE)
    add_textbox(slide, 0.45, y + 0.32, 5.8, 0.24, fdesc, font_size=9, color=C_LIGHT_GREY)
    y += 0.68

# Key efficiency features
add_textbox(slide, 6.7, 1.0, 6.3, 0.3, "KEY EFFICIENCY FEATURES", font_size=11, bold=True, color=C_ACCENT)

features = [
    ("NO_ERROR Short-Circuit",
     "The most impactful optimization in the entire pipeline. If _has_error_signals() returns False, the LLM is NEVER called and the pipeline terminates. This is the most common case in production — healthy traces — saving ~$0.10–$0.30 per pipeline run."),
    ("Two-Phase Error Detection",
     "Phase 1: explicit level check (O(1) dict lookup). Phase 2: keyword scan only for ambiguous levels, using a pre-compiled whole-word regex that avoids false positives from metric names like 'error_rate=0'."),
    ("Schema Built Once at Startup",
     "NormalizedIncident.model_json_schema() is called in __init__ and cached as self._response_schema. Every request reuses the same schema string — no repeated reflection."),
    ("Async Non-Blocking",
     "AsyncOpenAI + httpx.AsyncClient — all I/O is non-blocking. FastAPI handles concurrent requests without threads."),
]

y = 1.35
for title, desc in features:
    add_rect(slide, 6.7, y, 6.3, 1.32, C_BOX_BG2)
    add_rect(slide, 6.7, y, 0.07, 1.32, C_GREEN)
    add_textbox(slide, 6.85, y + 0.08, 6.1, 0.28, title, font_size=10.5, bold=True, color=C_GREEN)
    add_textbox(slide, 6.85, y + 0.4, 6.1, 0.82, desc, font_size=9.5, color=C_WHITE)
    y += 1.42


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Agent 2: Correlation
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "Agent 2 — Correlation Agent",
             "Purpose: Build a cross-system causal failure graph and decide where to route further analysis")

add_rect(slide, 0.3, 1.0, 12.7, 0.55, C_BOX_BG)
add_rect(slide, 0.3, 1.0, 0.08, 0.55, C_AMBER)
add_textbox(slide, 0.5, 1.06, 12.3, 0.42,
            "CRITICAL OUTPUT: analysis_target  —  the routing decision that controls which data sources ALL subsequent agents use",
            font_size=10, bold=True, color=C_AMBER)

# Left: Data fetching
add_textbox(slide, 0.3, 1.7, 6.0, 0.28, "DATA FETCHING STRATEGY", font_size=10, bold=True, color=C_ACCENT)

add_rect(slide, 0.3, 2.02, 5.9, 0.58, C_GREEN)
add_textbox(slide, 0.5, 2.09, 5.6, 0.22, "ALWAYS fetched", font_size=9.5, bold=True, color=C_NAVY)
add_textbox(slide, 0.5, 2.32, 5.6, 0.24, "Prometheus  →  6 PromQL queries  |  infrastructure baseline", font_size=9.5, color=C_NAVY)

add_rect(slide, 0.3, 2.65, 5.9, 0.58, C_ACCENT)
add_textbox(slide, 0.5, 2.72, 5.6, 0.22, "CONDITIONAL (if trace_id available)", font_size=9.5, bold=True, color=C_WHITE)
add_textbox(slide, 0.5, 2.95, 5.6, 0.24, "Langfuse  →  trace + up to 100 observations  |  AI layer view", font_size=9.5, color=C_WHITE)

add_rect(slide, 0.3, 3.28, 5.9, 0.45, C_BOX_BG)
add_textbox(slide, 0.5, 3.35, 5.6, 0.35,
            "All logs sorted chronologically before LLM — essential for causality (earliest failure = root candidate)",
            font_size=9.5, color=C_LIGHT_GREY, italic=True)

# Processing steps
add_textbox(slide, 0.3, 3.85, 6.0, 0.28, "PROCESSING STEPS", font_size=10, bold=True, color=C_ACCENT)
steps = [
    "Fetch Prometheus (always) + Langfuse (if trace_id)",
    "Sort all logs chronologically across both sources",
    "Build user message: logs grouped by source + normalization context",
    "Call GPT-4o with CorrelationResult schema",
    "Pydantic validate and return CorrelationResult + analysis_target",
]
y = 4.18
for i, s in enumerate(steps):
    add_rect(slide, 0.3, y, 0.35, 0.3, C_ACCENT)
    add_textbox(slide, 0.3, y + 0.03, 0.35, 0.24, str(i+1),
                font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.7, y + 0.03, 5.4, 0.28, s, font_size=9.5, color=C_WHITE)
    y += 0.38

# Right: Output + routing
add_textbox(slide, 6.7, 1.7, 6.3, 0.28, "OUTPUT: CorrelationResult", font_size=10, bold=True, color=C_ACCENT)

output_fields = [
    ("correlation_chain",    "list[str]  —  e.g. ['DNS failure → proxy timeout → gateway 502']"),
    ("peer_components",      "list[PeerComponent]  —  each with role + evidence"),
    ("timeline",             "list[TimelineEvent]  —  chronological failure events"),
    ("root_cause_candidate", "component + confidence (0–1) + reason"),
    ("analysis_target",      "Agent | InfraLogs | Unknown  ←  THE ROUTING DECISION"),
]
y = 2.02
for fname, fdesc in output_fields:
    bg = C_RED if fname == "analysis_target" else C_BOX_BG
    fc = C_AMBER if fname == "analysis_target" else C_AMBER
    add_rect(slide, 6.7, y, 6.3, 0.52, bg)
    add_textbox(slide, 6.85, y + 0.04, 2.2, 0.22, fname, font_size=9.5, bold=True, color=fc)
    add_textbox(slide, 6.85, y + 0.26, 6.1, 0.22, fdesc, font_size=9, color=C_WHITE)
    y += 0.57

# Routing table
add_textbox(slide, 6.7, 5.05, 6.3, 0.28, "ROUTING DECISION — analysis_target Values", font_size=10, bold=True, color=C_ACCENT)
table_row(slide, 6.7, 5.38, [2.0, 2.1, 2.2],
          ["analysis_target", "Agents 3 & 4 fetch", "Use when"],
          [C_DARK_BLUE, C_DARK_BLUE, C_DARK_BLUE], heights=0.3, bold_flags=[True, True, True])
routing_rows = [
    ("Agent",     "Langfuse only",           "Error is in AI agent logic / LLM calls"),
    ("InfraLogs", "Prometheus only",          "Error is in infrastructure / resources"),
    ("Unknown",   "Both Langfuse + Prometheus","Error source unclear — cross-check both"),
]
y = 5.68
bgs = [C_BOX_BG, C_BOX_BG2, C_BOX_BG]
for i, (target, fetch, when) in enumerate(routing_rows):
    table_row(slide, 6.7, y, [2.0, 2.1, 2.2], [target, fetch, when],
              [bgs[i], bgs[i], bgs[i]], heights=0.3)
    y += 0.3


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Agent 3: Error Analysis
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "Agent 3 — Error Analysis Agent",
             "Purpose: Deep-dive error extraction — every error gets a unique ID, category, severity, and evidence")

add_rect(slide, 0.3, 1.0, 12.7, 0.55, C_BOX_BG)
add_rect(slide, 0.3, 1.0, 0.08, 0.55, C_ACCENT)
add_textbox(slide, 0.5, 1.06, 12.3, 0.42,
            "STRICT SCOPE: Error identification and classification only. No root cause. No recommendations. No speculation.",
            font_size=10, bold=True, color=C_WHITE)

# Left: routing + steps
add_textbox(slide, 0.3, 1.7, 5.9, 0.28, "TARGETED DATA FETCHING", font_size=10, bold=True, color=C_ACCENT)

table_row(slide, 0.3, 2.02, [2.0, 1.95, 1.95],
          ["analysis_target", "Langfuse", "Prometheus"],
          [C_DARK_BLUE]*3, heights=0.3, bold_flags=[True, True, True])
fetch_rows = [
    ("Agent",     "✓  Yes (needs trace_id)", "✗  No"),
    ("InfraLogs", "✗  No",                   "✓  Yes"),
    ("Unknown",   "✓  Yes (needs trace_id)", "✓  Yes"),
]
y = 2.32
for i, (target, lf, pr) in enumerate(fetch_rows):
    bg = C_BOX_BG if i % 2 == 0 else C_BOX_BG2
    table_row(slide, 0.3, y, [2.0, 1.95, 1.95], [target, lf, pr], [bg]*3, heights=0.3)
    y += 0.3

add_textbox(slide, 0.3, 3.05, 5.9, 0.28, "PROCESSING STEPS", font_size=10, bold=True, color=C_ACCENT)
steps = [
    "Read analysis_target from CorrelationResult",
    "Fetch only the targeted data source(s)",
    "Merge + sort all logs chronologically",
    "Build prompt: correlation context + normalization + logs",
    "Call GPT-4o with ErrorAnalysisResult schema",
    "Pydantic validate — reject any missing error evidence",
    "Return ErrorAnalysisResult + rca_target (passthrough)",
]
y = 3.35
for i, s in enumerate(steps):
    add_rect(slide, 0.3, y, 0.32, 0.28, C_ACCENT)
    add_textbox(slide, 0.3, y + 0.02, 0.32, 0.24, str(i+1),
                font_size=9.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.67, y + 0.02, 5.4, 0.26, s, font_size=9.5, color=C_WHITE)
    y += 0.33

# Right: output
add_textbox(slide, 6.7, 1.7, 6.3, 0.28, "OUTPUT: ErrorAnalysisResult", font_size=10, bold=True, color=C_ACCENT)

add_rect(slide, 6.7, 2.02, 6.3, 0.28, C_DARK_BLUE)
add_textbox(slide, 6.85, 2.05, 6.1, 0.22, "Top-level fields", font_size=9, bold=True, color=C_LIGHT_GREY)
top_fields = [
    ("analysis_summary",       "str  (max 500 chars)"),
    ("analysis_target",        "Agent | InfraLogs | Unknown"),
    ("errors",                 "list[ErrorDetail]  — min 1 required"),
    ("error_patterns",         "list[ErrorPattern]  — requires 2+ occurrences"),
    ("error_impacts",          "list[ErrorImpact]  — per-component assessment"),
    ("error_propagation_path", "list[str]  — time-ordered propagation"),
    ("confidence",             "float  0.0 – 1.0"),
]
y = 2.3
for fname, ftype in top_fields:
    bg = C_BOX_BG if y % 0.6 > 0.3 else C_BOX_BG2
    add_rect(slide, 6.7, y, 6.3, 0.3, C_BOX_BG)
    add_textbox(slide, 6.85, y + 0.04, 2.5, 0.22, fname, font_size=9.5, bold=True, color=C_AMBER)
    add_textbox(slide, 9.4, y + 0.04, 3.5, 0.22, ftype, font_size=9.5, color=C_WHITE)
    y += 0.32

add_textbox(slide, 6.7, 4.55, 6.3, 0.28, "ErrorDetail Schema (per error)", font_size=10, bold=True, color=C_ACCENT)
add_rect(slide, 6.7, 4.85, 6.3, 1.85, C_BOX_BG2)
detail_fields = [
    ("error_id",      '"ERR-001", "ERR-002" ...  ← unique sequential ID'),
    ("category",      "llm_failure | timeout | dns_failure | configuration_error | ..."),
    ("severity",      "critical | high | medium | low | info"),
    ("component",     "exact service/component name from logs"),
    ("error_message", "exact error text from logs"),
    ("evidence",      "raw log line proving this error exists"),
    ("source",        "langfuse | prometheus"),
]
y = 4.9
for fname, fdesc in detail_fields:
    add_textbox(slide, 6.85, y, 1.6, 0.24, fname, font_size=9, bold=True, color=C_AMBER)
    add_textbox(slide, 8.5, y, 4.4, 0.24, fdesc, font_size=9, color=C_WHITE)
    y += 0.26

add_rect(slide, 6.7, 6.75, 6.3, 0.45, C_GREEN)
add_textbox(slide, 6.85, 6.8, 6.1, 0.35,
            "error_id cross-references: RCA links findings to specific errors by ID. Recommendation links solutions to specific errors by ID.",
            font_size=9, bold=True, color=C_NAVY)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Agent 4: RCA Overview
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "Agent 4 — RCA Agent",
             "Purpose: Determine definitively WHY errors occurred — root cause, causal chain, Five Whys")

add_rect(slide, 0.3, 1.0, 12.7, 0.55, C_BOX_BG)
add_rect(slide, 0.3, 1.0, 0.08, 0.55, C_RED)
add_textbox(slide, 0.5, 1.06, 12.3, 0.42,
            "Most analytically demanding agent. STRICT SCOPE: Root cause determination only. No recommendations. No remediation steps.",
            font_size=10, bold=True, color=C_WHITE)

# Left: steps
add_textbox(slide, 0.3, 1.7, 6.0, 0.28, "PROCESSING STEPS", font_size=10, bold=True, color=C_ACCENT)
steps = [
    ("1", "Read rca_target",        "Passed through from Error Analysis (= analysis_target)"),
    ("2", "Route data fetch",       "Agent→Langfuse only  |  InfraLogs→Prometheus only  |  Unknown→Both"),
    ("3", "Re-fetch fresh logs",    "Independent fetch from same sources — full context for causality"),
    ("4", "Build rich user message","All ErrorDetail + fresh logs + incident context"),
    ("5", "Build system prompt",    "Error categories, propagation path, Five Whys rules, RCAResult schema"),
    ("6", "Call GPT-4o",            "temperature=0.0  ·  Full RCAResult with 5-Why enforced"),
    ("7", "Pydantic validate",      "RCAResult.model_validate — rejects if five_why_analysis missing or incomplete"),
]
y = 2.02
for num, title, desc in steps:
    add_rect(slide, 0.3, y, 0.35, 0.56, C_ACCENT)
    add_textbox(slide, 0.3, y + 0.1, 0.35, 0.35, num,
                font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.65, y, 5.55, 0.56, C_BOX_BG)
    add_textbox(slide, 0.8, y + 0.04, 5.3, 0.24, title, font_size=10, bold=True, color=C_AMBER)
    add_textbox(slide, 0.8, y + 0.3, 5.3, 0.24, desc, font_size=9.5, color=C_LIGHT_GREY)
    y += 0.63

# Right: output fields
add_textbox(slide, 6.7, 1.7, 6.3, 0.28, "OUTPUT: RCAResult", font_size=10, bold=True, color=C_ACCENT)

rca_fields = [
    ("rca_summary",          "str  max 800 chars — executive summary"),
    ("root_cause",           "RootCause  (category, component, description, evidence, error_ids, confidence)"),
    ("causal_chain",         "list[CausalLink]  min 1  —  source→target with link_type + evidence"),
    ("contributing_factors", "list[ContributingFactor]  —  amplifiers that worsened the failure"),
    ("failure_timeline",     "list[FailureTimeline]  —  chronological events with is_root_cause flag"),
    ("blast_radius",         "list[str]  —  every component affected by the root cause"),
    ("five_why_analysis",    "FiveWhyAnalysis  —  EXACTLY 5 WhySteps + fundamental_root_cause  ← NEW"),
    ("confidence",           "float  0.0 – 1.0  —  tied to evidence strength"),
]
y = 2.02
for fname, fdesc in rca_fields:
    is_new = "NEW" in fdesc
    bg = C_RED if is_new else C_BOX_BG
    add_rect(slide, 6.7, y, 6.3, 0.56, bg)
    fc = C_AMBER if is_new else C_AMBER
    add_textbox(slide, 6.85, y + 0.04, 2.3, 0.24, fname, font_size=9.5, bold=True, color=fc)
    add_textbox(slide, 6.85, y + 0.3, 6.1, 0.22, fdesc, font_size=9, color=C_WHITE)
    y += 0.6

add_rect(slide, 6.7, 6.82, 6.3, 0.38, C_AMBER)
add_textbox(slide, 6.85, 6.87, 6.1, 0.28,
            "CausalLink types:  direct_cause  |  indirect_cause  |  trigger  |  amplifier",
            font_size=9.5, bold=True, color=C_NAVY)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Agent 4: Five Whys Analysis
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "Agent 4 — Five Whys Analysis",
             "Iterative technique: each answer becomes the subject of the next Why — drills from symptom to fundamental cause")

# Methodology
add_textbox(slide, 0.3, 1.0, 12.7, 0.28, "WhyStep Schema (×5 required)", font_size=10, bold=True, color=C_ACCENT)
schema_fields = [
    ("step",      "int 1–5",       "Position in the sequence"),
    ("question",  "str",           '"Why did [previous answer] occur?"'),
    ("answer",    "str",           "Explanation of the cause at this level"),
    ("evidence",  "str",           "Specific log line or metric supporting the answer"),
    ("component", "str",           "Implicated service or system element"),
]
table_row(slide, 0.3, 1.32, [1.4, 1.4, 9.2],
          ["Field", "Type", "Description"],
          [C_DARK_BLUE]*3, heights=0.28, bold_flags=[True]*3)
y = 1.6
for fname, ftype, fdesc in schema_fields:
    bg = C_BOX_BG if y % 0.6 > 0.3 else C_BOX_BG2
    table_row(slide, 0.3, y, [1.4, 1.4, 9.2], [fname, ftype, fdesc], [bg]*3, heights=0.28)
    y += 0.28

# Example
add_textbox(slide, 0.3, 2.9, 12.7, 0.28,
            "EXAMPLE  —  LLM Access Disabled in 'sample-agent' service", font_size=10, bold=True, color=C_ACCENT)

example_whys = [
    ("Problem", "sample-agent fails to process requests — LLM access is disabled",
     "LLM access is disabled (demo error mode) log entry"),
    ("Why 1",   "LLM access is disabled in the service",
     "Log: LLM access is disabled (demo error mode). Click Enable LLM Access..."),
    ("Why 2",   "Service is configured to run in demo error mode",
     "Error message explicitly references 'demo error mode' as the cause"),
    ("Why 3",   "Demo error mode config flag was toggled (manually via UI or deployment)",
     "Error instructs user to click 'Enable LLM Access' in chat UI — UI toggle controls this"),
    ("Why 4",   "No pre-flight validation checks LLM access state before accepting requests",
     "Service accepted the request and only failed at the LLM call stage — no early-rejection log"),
    ("Why 5",   "Demo error mode was added for testing without a corresponding readiness gate",
     "Inferred from absence of startup/readiness logs rejecting demo mode (limited visibility)"),
]

colors_ex = [C_DARK_BLUE, C_BOX_BG, C_BOX_BG, C_BOX_BG, C_BOX_BG, C_RED]
label_colors = [C_LIGHT_GREY, C_WHITE, C_WHITE, C_WHITE, C_WHITE, C_AMBER]

y = 3.22
for i, (label, answer, evidence) in enumerate(example_whys):
    add_rect(slide, 0.3, y, 1.2, 0.55, colors_ex[i])
    add_textbox(slide, 0.3, y + 0.12, 1.2, 0.3, label,
                font_size=10, bold=True, color=label_colors[i], align=PP_ALIGN.CENTER)
    add_rect(slide, 1.5, y, 11.5, 0.55, C_BOX_BG if i % 2 == 0 else C_BOX_BG2)
    add_textbox(slide, 1.65, y + 0.02, 11.2, 0.24, answer, font_size=9.5, bold=(i==5), color=C_WHITE)
    add_textbox(slide, 1.65, y + 0.3, 11.2, 0.22, f"Evidence: {evidence}", font_size=8.5, color=C_LIGHT_GREY, italic=True)
    y += 0.6

add_rect(slide, 0.3, 6.88, 12.7, 0.38, C_GREEN)
add_textbox(slide, 0.5, 6.93, 12.3, 0.28,
            "Fundamental Root Cause:  Demo error mode feature lacks a readiness gate — configuration toggle was never paired with an enforcement mechanism",
            font_size=9.5, bold=True, color=C_NAVY)

# Enforcement note
add_rect(slide, 0.3, y + 0.05, 12.7, 0.0, C_BOX_BG)  # spacer


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Agent 5: Recommendation
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "Agent 5 — Recommendation Agent",
             "Purpose: Synthesise Error Analysis + RCA into 1–4 ranked, actionable solutions")

add_rect(slide, 0.3, 1.0, 12.7, 0.52, C_GREEN)
add_textbox(slide, 0.5, 1.07, 12.3, 0.38,
            "ONLY agent that fetches ZERO external data. Purely a synthesis layer — no Langfuse, no Prometheus, no I/O wait. Fastest step in pipeline.",
            font_size=10, bold=True, color=C_NAVY)

# Left
add_textbox(slide, 0.3, 1.68, 6.0, 0.28, "PROCESSING STEPS", font_size=10, bold=True, color=C_ACCENT)
steps = [
    ("1", "Build rich user message",  "RCA summary + causal chain + contributing factors + blast radius + all errors + patterns + impacts"),
    ("2", "Build system prompt",      "13 context fields injected from RCA (8) + Error Analysis (5) + ranking rules + solution schema"),
    ("3", "Call GPT-4o",              "temperature=0.0  ·  1–4 solutions only — anti-padding rule enforced in prompt"),
    ("4", "Pydantic validate",        "model_validator checks ranks are sequential from 1 to N — rejects gaps or duplicates"),
]
y = 2.0
for num, title, desc in steps:
    add_rect(slide, 0.3, y, 0.35, 0.88, C_ACCENT)
    add_textbox(slide, 0.3, y + 0.22, 0.35, 0.4, num,
                font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.65, y, 5.55, 0.88, C_BOX_BG)
    add_textbox(slide, 0.8, y + 0.08, 5.3, 0.28, title, font_size=10, bold=True, color=C_AMBER)
    add_textbox(slide, 0.8, y + 0.42, 5.3, 0.4, desc, font_size=9.5, color=C_LIGHT_GREY)
    y += 0.98

# Ranking rules
add_textbox(slide, 0.3, 5.65, 6.0, 0.28, "RANKING RULES", font_size=10, bold=True, color=C_ACCENT)
ranking = [
    ("Rank 1", C_RED,    "MUST directly fix the root cause — addresses_root_cause=True"),
    ("Rank 2", C_AMBER,  "Most critical secondary concern or prevents propagation"),
    ("Rank 3", C_ACCENT, "Contributing factor or resilience improvement"),
    ("Rank 4", C_GREEN,  "Preventive measure for future recurrence (only if genuinely useful)"),
]
y = 5.95
for rank, color, rule in ranking:
    add_rect(slide, 0.3, y, 1.0, 0.3, color)
    add_textbox(slide, 0.3, y + 0.04, 1.0, 0.22, rank,
                font_size=9.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.35, y + 0.04, 4.85, 0.24, rule, font_size=9.5, color=C_WHITE)
    y += 0.35

# Right: Solution schema
add_textbox(slide, 6.7, 1.68, 6.3, 0.28, "OUTPUT: Solution Schema", font_size=10, bold=True, color=C_ACCENT)
solution_fields = [
    ("rank",                "int 1–4  — unique, sequential (Pydantic validated)"),
    ("title",               "str  max 120 chars  — short actionable title"),
    ("description",         "str  — detailed action + why it addresses root cause"),
    ("category",            "config_change | code_fix | infrastructure | scaling | retry_logic | fallback | monitoring | ..."),
    ("effort",              "quick_fix | low | medium | high"),
    ("addresses_root_cause","bool  — True only for solutions that fix the root cause directly"),
    ("affected_components", "list[str]  — components this solution targets"),
    ("expected_outcome",    "str  — what improvement is expected after implementation"),
    ("error_ids",           "list[str]  — links back to ErrorAnalysis error IDs (ERR-001, ...)"),
]
y = 2.0
for fname, fdesc in solution_fields:
    add_rect(slide, 6.7, y, 6.3, 0.52, C_BOX_BG if y % 1.0 > 0.5 else C_BOX_BG2)
    add_textbox(slide, 6.85, y + 0.03, 2.4, 0.22, fname, font_size=9.5, bold=True, color=C_AMBER)
    add_textbox(slide, 6.85, y + 0.28, 6.1, 0.2, fdesc, font_size=9, color=C_WHITE)
    y += 0.56


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Pipeline Data Flow
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "End-to-End Pipeline Data Flow",
             "What each agent receives, what it fetches externally, and what it passes to the next agent")

flow = [
    ("Normalization",   "timestamp\ntrace_id?\nagent_name",
     "Langfuse (if trace_id)\nPrometheus (if no trace_id)",
     "NormalizedIncident\nerror_type, signals,\nentities, confidence"),
    ("Correlation",     "NormalizedIncident\ntrace_id?\nagent_name",
     "Prometheus (always)\nLangfuse (if trace_id)",
     "CorrelationResult\ncausal chain, timeline,\nanalysis_target  ←KEY"),
    ("Error Analysis",  "CorrelationResult\nNormalizedIncident\ntrace_id? + agent_name",
     "Langfuse (if Agent/Unk)\nPrometheus (if Infra/Unk)",
     "ErrorAnalysisResult\nerrors with IDs,\npatterns, rca_target"),
    ("RCA",             "ErrorAnalysisResult\nNormalizedIncident\nrca_target + trace_id?",
     "Langfuse (if Agent/Unk)\nPrometheus (if Infra/Unk)",
     "RCAResult\nroot_cause, causal_chain\nfive_why_analysis"),
    ("Recommendation",  "ErrorAnalysisResult\nRCAResult\nagent_name",
     "NONE — synthesis only",
     "RecommendationResult\n1–4 ranked solutions\neffort + category"),
]

# Header
table_row(slide, 0.3, 1.0, [2.1, 3.2, 3.4, 4.0],
          ["Agent", "Receives from prev agent", "Fetches externally", "Passes to next agent"],
          [C_DARK_BLUE]*4, heights=0.35, bold_flags=[True]*4)

y = 1.35
row_colors = [C_BOX_BG, C_BOX_BG2, C_BOX_BG, C_BOX_BG2, C_BOX_BG]
for i, (agent, receives, fetches, passes) in enumerate(flow):
    h = 1.05
    fetch_color = C_BOX_BG2 if "NONE" in fetches else row_colors[i]
    fc_text = C_GREEN if "NONE" in fetches else C_WHITE

    add_rect(slide, 0.3, y, 2.1, h, C_ACCENT)
    add_textbox(slide, 0.3, y + 0.3, 2.1, 0.4, agent,
                font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, 2.4, y, 3.2, h, row_colors[i])
    add_textbox(slide, 2.5, y + 0.08, 3.0, h - 0.15, receives, font_size=9, color=C_WHITE)

    add_rect(slide, 5.6, y, 3.4, h, fetch_color)
    add_textbox(slide, 5.7, y + 0.08, 3.2, h - 0.15, fetches, font_size=9, color=fc_text)

    add_rect(slide, 9.0, y, 4.1, h, row_colors[i])
    add_textbox(slide, 9.1, y + 0.08, 3.9, h - 0.15, passes, font_size=9, color=C_WHITE)

    y += h + 0.04

add_textbox(slide, 0.3, 6.92, 12.7, 0.28,
            "External data is fetched independently by each agent that needs it — no shared log cache. Re-fetching ensures each agent gets full, unfiltered evidence for its specific analysis task.",
            font_size=9, color=C_LIGHT_GREY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Efficiency & Design Decisions
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "Efficiency & Key Design Decisions")

decisions = [
    (C_GREEN,  "NO_ERROR Short-Circuit",
     "Pre-LLM scan in Normalization: if no error signals detected → return immediately, skip all 5 agents. Most common production case (healthy traces). Saves ~$0.10–$0.30 per run."),
    (C_ACCENT, "Schema-Driven LLM Contracts",
     "Every Pydantic model auto-generates its JSON schema via model_json_schema(), injected into the system prompt. Adding a new field automatically updates the LLM contract — no manual JSON maintenance."),
    (C_AMBER,  "Error ID Cross-Referencing",
     "Error Analysis assigns ERR-001, ERR-002, ... to each distinct error. RCA references these IDs in causal links and root_cause.error_ids. Recommendation references them in each solution.error_ids. Lightweight linking without re-embedding full error objects."),
    (C_RED,    "Routing Decisions Flow Forward",
     "Correlation sets analysis_target once. Error Analysis passes it as rca_target. RCA reads rca_target. Each downstream agent inherits the routing decision without re-analysing the situation."),
    (C_GREEN,  "Graceful Degradation",
     "Every external source failure (Langfuse, Prometheus) produces a WARN placeholder entry. The LLM is explicitly told data was missing and instructed to set lower confidence. No agent returns an error; the pipeline always completes."),
    (C_ACCENT, "Five Whys Schema Enforcement",
     "FiveWhyAnalysis.whys has min_length=5, max_length=5 in Pydantic. The LLM cannot return fewer or more than 5 steps. The prompt defines the exact chaining methodology. Evidence is required at every step."),
]

y = 1.05
for i, (color, title, desc) in enumerate(decisions):
    col = i % 2
    row = i // 2
    lx = 0.3 + col * 6.55
    ly = y + row * 1.85
    add_rect(slide, lx, ly, 6.3, 1.72, C_BOX_BG)
    add_rect(slide, lx, ly, 0.08, 1.72, color)
    add_textbox(slide, lx + 0.22, ly + 0.1, 5.95, 0.3, title, font_size=11, bold=True, color=color)
    add_textbox(slide, lx + 0.22, ly + 0.45, 5.95, 1.15, desc, font_size=9.5, color=C_WHITE)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Summary Quick Reference
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_NAVY)
slide_header(slide, "Summary — Quick Reference", "All five agents at a glance")

# Summary table
cols = [2.2, 2.0, 2.5, 2.5, 3.8]
table_row(slide, 0.15, 1.0, cols,
          ["Agent", "Data Sources", "Key Output", "Critical Field", "Scope Boundary"],
          [C_DARK_BLUE]*5, heights=0.35, bold_flags=[True]*5)

summary_rows = [
    ("1  Normalization",   "Langfuse OR Prometheus\n(not both — mutual exclusive)",
     "NormalizedIncident", "error_type",      "Classify error only — no causality"),
    ("2  Correlation",     "Prometheus (always)\n+ Langfuse (if trace_id)",
     "CorrelationResult",  "analysis_target", "Causal graph only — no per-error detail"),
    ("3  Error Analysis",  "Langfuse OR Prometheus\n(based on analysis_target)",
     "ErrorAnalysisResult","error_ids (ERR-N)","Error extraction only — no root cause"),
    ("4  RCA",             "Langfuse OR Prometheus\n(based on rca_target)",
     "RCAResult",          "five_why_analysis","Root cause only — no recommendations"),
    ("5  Recommendation",  "NONE\n(synthesis only)",
     "RecommendationResult","rank + error_ids","Solutions only — no diagnosis"),
]
y = 1.35
row_bgs = [C_BOX_BG, C_BOX_BG2, C_BOX_BG, C_BOX_BG2, C_BOX_BG]
for i, (agent, sources, output, key_field, boundary) in enumerate(summary_rows):
    h = 0.88
    table_row(slide, 0.15, y, cols,
              [agent, sources, output, key_field, boundary],
              [row_bgs[i]]*5, heights=h)
    y += h + 0.02

# Bottom row: stats
stats = [
    ("5", "Agents in pipeline"),
    ("2", "External data sources\n(Langfuse + Prometheus)"),
    ("6", "PromQL queries per\nPrometheus fetch"),
    ("0.0", "LLM temperature\n(deterministic outputs)"),
    ("5", "WhyStep responses\nin Five Whys analysis"),
    ("1–4", "Ranked solutions from\nRecommendation Agent"),
]
x = 0.15
for stat, label in stats:
    add_rect(slide, x, 5.95, 2.15, 1.3, C_BOX_BG)
    add_textbox(slide, x, 6.02, 2.15, 0.65, stat,
                font_size=30, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, 6.67, 2.15, 0.5, label,
                font_size=8.5, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER)
    x += 2.2

# Footer
add_textbox(slide, 0.3, 7.18, 12.7, 0.28,
            "Investigation Pipeline  ·  Multi-Agent Observability System  ·  Prodapt AI Team  ·  2026",
            font_size=9, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)


# ── Save ───────────────────────────────────────────────────────────────────
output_path = r"C:\Users\vyanktesh.l\Documents\Invastigate_flow_with_Poller\doc_for_reference\Agent_Pipeline_Documentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
